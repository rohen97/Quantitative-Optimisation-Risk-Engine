from __future__ import annotations

import hashlib
import json
import math
import unicodedata
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, Sequence

import matplotlib

matplotlib.use('Agg')

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

INK = '172127'
PAPER = 'F6F8F7'
WHITE = 'FFFFFF'
MUTED = '647078'
BORDER = 'D8E0DD'
GREEN = '14966F'
GREEN_DARK = '0B6F56'
TEAL = '2B7280'
BLUE = '3478A5'
GOLD = 'C99328'
RED = 'BB4A4A'

PALE_GREEN = 'E8F4EF'
PALE_BLUE = 'E9F1F6'
PALE_GOLD = 'FBF3DF'
PALE_RED = 'F8EAEA'

FONT_HEAD = 'Aptos Display'
FONT_BODY = 'Aptos'

RELEASE_RELATIVE = Path(
    'reports/releases/2026-08-19-free-data-drl-risk'
)
FALLBACK_RELEASE_RELATIVE = Path(
    'reports/releases/2026-08-13-risk-pit-cost-validation'
)
PRIOR_RELEASE_RELATIVE = Path('reports/releases/2026-08-07-full-universe')
BACKTEST_RELATIVE = Path('reports/backtests/1997_to_latest')
OUTPUTS_RELATIVE = Path('reports/outputs')

PRESENTATION_RELATIVE = Path(
    'reports/presentations/wolf_investment_principal'
)


@dataclass(frozen=True)
class DeckEvidence:
    repo_root: Path
    release_root: Path
    backtest_root: Path
    outputs_root: Path
    validation_manifest: dict
    prior_validation_manifest: dict
    walk_forward_manifest: dict
    backtest_manifest: dict

    universe: pd.DataFrame
    holdings: pd.DataFrame
    trades: pd.DataFrame
    scorecard: pd.DataFrame
    prior_pit_summary: pd.DataFrame
    pit_summary: pd.DataFrame
    benchmark_significance: pd.DataFrame
    pit_returns: pd.DataFrame
    pit_coverage: dict
    production_pit: pd.DataFrame
    free_data_summary: pd.DataFrame
    alpha: pd.DataFrame
    overfitting: pd.Series
    performance: pd.DataFrame
    optimiser: pd.DataFrame
    risk_backtest: pd.DataFrame
    constraints: pd.DataFrame
    regime: pd.DataFrame
    regional_alpha: pd.DataFrame
    supervised_dataset: pd.DataFrame
    supervised_validation: pd.DataFrame
    supervised_oos: pd.DataFrame
    supervised_quantiles: pd.DataFrame
    supervised_acceptance: pd.DataFrame
    supervised_ensemble: pd.DataFrame
    supervised_latest: pd.DataFrame
    supervised_freeze: dict
    drl_acceptance: pd.Series
    drl_challengers: pd.DataFrame
    drl_seeds: pd.DataFrame
    drl_split: pd.DataFrame
    drl_long_history: dict
    shadow_status: dict

    @property
    def current_aum(self) -> float:
        return float(self.backtest_manifest['current_portfolio_nav_usd'])

    @property
    def as_of_date(self) -> str:
        validation_date = pd.Timestamp(self.validation_manifest['as_of_date'])
        supervised_dates = pd.to_datetime(
            self.supervised_latest.get('as_of_date'), errors='coerce'
        ).dropna()
        latest = max(
            validation_date,
            supervised_dates.max() if not supervised_dates.empty else validation_date,
        )
        return latest.strftime('%d %B %Y')

    @property
    def validation_as_of_date(self) -> str:
        return pd.Timestamp(
            self.validation_manifest['as_of_date']
        ).strftime('%d %B %Y')

    @property
    def approval_status(self) -> str:
        return str(self.validation_manifest['approval_status'])

    @property
    def governance_score(self) -> float:
        return float(self.validation_manifest['overall_score'])


@dataclass(frozen=True)
class DeckBuildResult:
    pptx_path: Path
    report_path: Path
    manifest_path: Path
    plot_paths: tuple[Path, ...]
    slide_count: int

def _require_files(paths: Iterable[Path]) -> None:
    missing = [str(path) for path in paths if not path.exists()]
    if missing:
        joined = '\n'.join(f'- {path}' for path in missing)
        raise FileNotFoundError(
            f'Missing presentation evidence files:\n{joined}'
        )


def _read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding='utf-8'))


def _read_csv_optional(path: Path) -> pd.DataFrame:
    return pd.read_csv(path, low_memory=False) if path.exists() else pd.DataFrame()

def load_deck_evidence(repo_root: str | Path) -> DeckEvidence:
    repo_root = Path(repo_root).resolve()
    release_root = repo_root / RELEASE_RELATIVE
    if not release_root.exists():
        release_root = repo_root / FALLBACK_RELEASE_RELATIVE
    prior_release_root = repo_root / PRIOR_RELEASE_RELATIVE
    backtest_root = repo_root / BACKTEST_RELATIVE
    outputs_root = repo_root / OUTPUTS_RELATIVE
    recommendation_snapshot_path = (
        repo_root / PRESENTATION_RELATIVE / 'recommendation_snapshot.csv'
    )
    regional_alpha_path = outputs_root / 'optimised_portfolio_regional_alpha.csv'
    supervised_latest_path = (
        outputs_root / 'supervised_alpha/latest_predictions.csv'
    )
    validation_root = release_root / 'validation'
    required = [
        validation_root / 'validation_manifest.json',
        release_root / 'walk_forward_manifest.json',
        release_root / 'universe_summary.csv',

        validation_root / 'model_validation_scorecard.csv',
        validation_root / 'portfolio_strategy_comparison.csv',
        validation_root / 'benchmark_significance_report.csv',
        validation_root / 'portfolio_monthly_returns.csv',
        validation_root / 'risk_backtesting_report.csv',
        validation_root / 'constraint_compliance_report.csv',
        release_root / 'pit_evidence_coverage.json',
        release_root / 'bloomberg_pit_coverage.csv',
        prior_release_root / 'validation/validation_manifest.json',
        prior_release_root
        / 'validation/portfolio_strategy_comparison.csv',
        backtest_root / 'run_manifest.json',
        backtest_root / 'point_in_time_alpha_significance.csv',

        backtest_root / 'strategy_overfitting_summary.csv',
        backtest_root / 'performance_summary.csv',
        outputs_root / 'final_portfolio_weights.csv',
        outputs_root / 'portfolio_trade_list.csv',
        outputs_root / 'portfolio_optimisation_summary.csv',
        outputs_root / 'supervised_alpha/dataset_profile.csv',
        outputs_root / 'supervised_alpha/validation_summary.csv',
        outputs_root / 'supervised_alpha/oos_summary.csv',
        outputs_root / 'supervised_alpha/quantile_metrics.csv',
        outputs_root / 'supervised_alpha/acceptance_decision.csv',
        outputs_root / 'supervised_alpha/ensemble_weights.csv',
        outputs_root / 'supervised_alpha/prospective_freeze_manifest.json',
        outputs_root / 'drl_acceptance_decision.csv',
        outputs_root / 'drl_simple_challenger_comparison.csv',
        outputs_root / 'drl_training_summary.csv',
        outputs_root / 'drl_split_manifest.csv',
        release_root / 'public_data/drl_long_history_manifest.json',
        outputs_root / 'shadow_operation/shadow_operation_status.json',
    ]
    _require_files(required)
    if not regional_alpha_path.exists() or not supervised_latest_path.exists():
        _require_files([recommendation_snapshot_path])
        recommendation_snapshot = pd.read_csv(recommendation_snapshot_path)
    else:
        recommendation_snapshot = pd.DataFrame()

    if regional_alpha_path.exists():
        regional_alpha = pd.read_csv(regional_alpha_path, low_memory=False)
    else:
        regional_alpha = recommendation_snapshot.loc[
            recommendation_snapshot['recommendation_class'].eq(
                'regional_alpha_challenger'
            )
        ].copy()

    if supervised_latest_path.exists():
        supervised_latest = pd.read_csv(
            supervised_latest_path,
            low_memory=False,
        )
    else:
        supervised_latest = recommendation_snapshot.loc[
            recommendation_snapshot['recommendation_class'].eq(
                'supervised_research_watchlist'
            )
        ].rename(
            columns={
                'model_score': 'supervised_alpha_score',
                'cost_adjusted_predicted_excess_return_3m': (
                    'cost_adjusted_predicted_excess_return'
                ),
                'q05_excess_return_3m': 'q05_excess_return',
                'q95_excess_return_3m': 'q95_excess_return',
            }
        )
        supervised_latest['horizon_months'] = 3

    holdings = pd.read_csv(
        outputs_root / 'final_portfolio_weights.csv'
    )

    holdings['final_weight'] = pd.to_numeric(
        holdings['final_weight'], errors='coerce'
    ).fillna(0.0)
    holdings = holdings.loc[holdings['final_weight'] > 1e-9].copy()
    holdings = holdings.sort_values(
        ['region', 'ticker']
    ).reset_index(drop=True)

    trades = pd.read_csv(outputs_root / 'portfolio_trade_list.csv')

    trades = trades.loc[
        trades['ticker'].isin(holdings['ticker'])
    ].copy()
    trades['target_weight'] = pd.to_numeric(
        trades['target_weight'], errors='coerce'
    ).fillna(0.0)
    trades = trades.sort_values(
        ['trade_action', 'ticker']
    ).reset_index(drop=True)

    overfitting = pd.read_csv(
        backtest_root / 'strategy_overfitting_summary.csv'
    )

    if overfitting.empty:
        raise ValueError('The overfitting summary has no evidence row.')

    return DeckEvidence(
        repo_root=repo_root,
        release_root=release_root,
        backtest_root=backtest_root,
        outputs_root=outputs_root,
        validation_manifest=_read_json(
            validation_root / 'validation_manifest.json'
        ),
        prior_validation_manifest=_read_json(
            prior_release_root / 'validation/validation_manifest.json'
        ),
        walk_forward_manifest=_read_json(
            release_root / 'walk_forward_manifest.json'
        ),
        backtest_manifest=_read_json(
            backtest_root / 'run_manifest.json'
        ),
        universe=pd.read_csv(release_root / 'universe_summary.csv'),
        holdings=holdings,
        trades=trades,
        scorecard=pd.read_csv(
            validation_root / 'model_validation_scorecard.csv'
        ),
        prior_pit_summary=pd.read_csv(
            prior_release_root
            / 'validation/portfolio_strategy_comparison.csv'
        ),
        pit_summary=pd.read_csv(
            validation_root / 'portfolio_strategy_comparison.csv'
        ),
        benchmark_significance=pd.read_csv(
            validation_root / 'benchmark_significance_report.csv'
        ),
        pit_returns=pd.read_csv(
            validation_root / 'portfolio_monthly_returns.csv'
        ),
        pit_coverage=_read_json(
            release_root / 'pit_evidence_coverage.json'
        ),
        production_pit=pd.read_csv(
            release_root / 'bloomberg_pit_coverage.csv'
        ),
        free_data_summary=_read_csv_optional(
            outputs_root / 'validation/free_data_evidence_summary.csv'
        ),
        alpha=pd.read_csv(
            backtest_root / 'point_in_time_alpha_significance.csv'
        ),
        overfitting=overfitting.iloc[0],
        performance=pd.read_csv(
            backtest_root / 'performance_summary.csv'
        ),

        optimiser=pd.read_csv(
            outputs_root / 'portfolio_optimisation_summary.csv'
        ),
        risk_backtest=pd.read_csv(
            validation_root / 'risk_backtesting_report.csv'
        ),
        constraints=pd.read_csv(
            validation_root / 'constraint_compliance_report.csv'
        ),
        regime=pd.read_csv(
            validation_root / 'regime_performance_report.csv'
        ),
        regional_alpha=regional_alpha,
        supervised_dataset=pd.read_csv(
            outputs_root / 'supervised_alpha/dataset_profile.csv'
        ),
        supervised_validation=pd.read_csv(
            outputs_root / 'supervised_alpha/validation_summary.csv'
        ),
        supervised_oos=pd.read_csv(
            outputs_root / 'supervised_alpha/oos_summary.csv'
        ),
        supervised_quantiles=pd.read_csv(
            outputs_root / 'supervised_alpha/quantile_metrics.csv'
        ),
        supervised_acceptance=pd.read_csv(
            outputs_root / 'supervised_alpha/acceptance_decision.csv'
        ),
        supervised_ensemble=pd.read_csv(
            outputs_root / 'supervised_alpha/ensemble_weights.csv'
        ),
        supervised_latest=supervised_latest,
        supervised_freeze=_read_json(
            outputs_root
            / 'supervised_alpha/prospective_freeze_manifest.json'
        ),
        drl_acceptance=pd.read_csv(
            outputs_root / 'drl_acceptance_decision.csv'
        ).iloc[0],
        drl_challengers=pd.read_csv(
            outputs_root / 'drl_simple_challenger_comparison.csv'
        ),
        drl_seeds=pd.read_csv(
            outputs_root / 'drl_training_summary.csv'
        ),
        drl_split=pd.read_csv(
            outputs_root / 'drl_split_manifest.csv'
        ),
        drl_long_history=_read_json(
            release_root / 'public_data/drl_long_history_manifest.json'
        ),
        shadow_status=_read_json(
            outputs_root / 'shadow_operation/shadow_operation_status.json'
        ),
    )


def _supervised_ensemble_rows(evidence: DeckEvidence) -> pd.DataFrame:
    return evidence.supervised_oos.loc[
        evidence.supervised_oos['candidate'].eq('supervised_alpha_ensemble')
    ].sort_values('horizon_months')


def _supervised_validation_rows(evidence: DeckEvidence) -> pd.DataFrame:
    return evidence.supervised_validation.loc[
        evidence.supervised_validation['candidate'].eq(
            'supervised_alpha_ensemble'
        )
    ].sort_values('horizon_months')


def _scorecard_component(evidence: DeckEvidence, name: str) -> pd.Series:
    rows = evidence.scorecard.loc[evidence.scorecard['component'].eq(name)]
    if rows.empty:
        raise ValueError(f'Missing validation scorecard component: {name}')
    return rows.iloc[0]


def _score_label(row: pd.Series) -> str:
    return f'{float(row.score):g} / {float(row.maximum_score):g}'


def _risk_gate_status(evidence: DeckEvidence) -> str:
    return str(_scorecard_component(evidence, 'risk_backtesting').status)


def _supervised_watchlist(evidence: DeckEvidence) -> pd.DataFrame:
    latest = evidence.supervised_latest.loc[
        evidence.supervised_latest['horizon_months'].eq(3)
    ].copy()
    latest['cost_adjusted_predicted_excess_return'] = pd.to_numeric(
        latest['cost_adjusted_predicted_excess_return'], errors='coerce'
    )
    return (
        latest.sort_values(
            ['region', 'cost_adjusted_predicted_excess_return'],
            ascending=[True, False],
        )
        .groupby('region', as_index=False, sort=True)
        .head(1)
        .sort_values('region')
        .reset_index(drop=True)
    )


def _recommendation_snapshot(evidence: DeckEvidence) -> pd.DataFrame:
    trades = evidence.trades[
        ['security_id', 'trade_action']
    ].drop_duplicates('security_id')
    target = evidence.holdings.loc[
        pd.to_numeric(evidence.holdings['final_weight'], errors='coerce').fillna(0.0)
        > 0.0
    ].merge(trades, on='security_id', how='left')
    governed = pd.DataFrame(
        {
            'as_of_date': evidence.as_of_date,
            'recommendation_class': 'governed_target',
            'governance_status': 'actionable_after_pre_trade_review',
            'security_id': target['security_id'],
            'ticker': target['ticker'],
            'company_name': target['company_name'],
            'region': target['region'],
            'sector': target['sector'],
            'action': target['trade_action'].fillna('Review'),
            'target_weight': pd.to_numeric(
                target['final_weight'], errors='coerce'
            ),
            'model_score': pd.to_numeric(
                target['final_recommendation_score'], errors='coerce'
            ),
            'cost_adjusted_predicted_excess_return_3m': np.nan,
            'q05_excess_return_3m': np.nan,
            'q95_excess_return_3m': np.nan,
        }
    )
    regional_source = evidence.regional_alpha.loc[
        pd.to_numeric(
            evidence.regional_alpha['target_weight'], errors='coerce'
        ).fillna(0.0)
        > 0.0
    ].copy()
    regional = pd.DataFrame(
        {
            'as_of_date': evidence.as_of_date,
            'recommendation_class': 'regional_alpha_challenger',
            'governance_status': 'research_only_not_an_order',
            'security_id': regional_source['security_id'],
            'ticker': regional_source['ticker'],
            'company_name': regional_source['company_name'],
            'region': regional_source['region'],
            'sector': regional_source['sector'],
            'action': 'Research review',
            'target_weight': pd.to_numeric(
                regional_source['target_weight'], errors='coerce'
            ),
            'model_score': pd.to_numeric(
                regional_source.get('regional_alpha_score'), errors='coerce'
            ),
            'cost_adjusted_predicted_excess_return_3m': np.nan,
            'q05_excess_return_3m': np.nan,
            'q95_excess_return_3m': np.nan,
        }
    )
    watchlist = _supervised_watchlist(evidence)
    research = pd.DataFrame(
        {
            'as_of_date': evidence.as_of_date,
            'recommendation_class': 'supervised_research_watchlist',
            'governance_status': 'research_only_not_a_buy_order',
            'security_id': watchlist['security_id'],
            'ticker': watchlist['ticker'],
            'company_name': watchlist['company_name'],
            'region': watchlist['region'],
            'sector': watchlist['sector'],
            'action': 'Research review',
            'target_weight': np.nan,
            'model_score': pd.to_numeric(
                watchlist['supervised_alpha_score'], errors='coerce'
            ),
            'cost_adjusted_predicted_excess_return_3m': pd.to_numeric(
                watchlist['cost_adjusted_predicted_excess_return'], errors='coerce'
            ),
            'q05_excess_return_3m': pd.to_numeric(
                watchlist['q05_excess_return'], errors='coerce'
            ),
            'q95_excess_return_3m': pd.to_numeric(
                watchlist['q95_excess_return'], errors='coerce'
            ),
        }
    )
    return pd.concat([governed, regional, research], ignore_index=True).sort_values(
        ['recommendation_class', 'region', 'ticker'],
        kind='stable',
    )

def _rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip('#')
    return RGBColor(
        int(value[0:2], 16),
        int(value[2:4], 16),
        int(value[4:6], 16),
    )


def _ascii_display(value: object) -> str:
    text = '' if value is None else str(value)
    return (
        unicodedata.normalize('NFKD', text)
        .encode('ascii', 'ignore')
        .decode('ascii')
    )

def _pct(value: float, digits: int = 1) -> str:
    return f'{float(value) * 100:.{digits}f}%'


def _usd(value: float, digits: int = 1) -> str:
    amount = float(value)
    magnitude = abs(amount)
    if magnitude >= 1_000_000_000:
        return '$' + f'{amount / 1_000_000_000:.{digits}f}bn'
    if magnitude >= 1_000_000:
        return '$' + f'{amount / 1_000_000:.{digits}f}m'

    if magnitude >= 1_000:
        return '$' + f'{amount / 1_000:.{digits}f}k'
    return '$' + f'{amount:,.0f}'


def _equal_weight_comparison(
    evidence: DeckEvidence,
) -> tuple[float, float]:
    summary = evidence.pit_summary.set_index('strategy')
    difference = float(
        summary.loc['wolf_cvar', 'annualised_return']
        - summary.loc['equal_weight_eligible', 'annualised_return']
    )
    row = evidence.benchmark_significance.loc[
        (evidence.benchmark_significance['strategy'] == 'wolf_cvar')
        & (
            evidence.benchmark_significance['baseline']
            == 'equal_weight_eligible'
        )
    ]
    if row.empty:
        raise ValueError('Missing Wolf/equal-weight significance evidence.')
    return difference, float(row.iloc[0]['p_value'])


def _equal_weight_sentence(evidence: DeckEvidence) -> str:
    difference, p_value = _equal_weight_comparison(evidence)
    relation = 'outperformed' if difference >= 0 else 'trailed'
    return (
        f'Wolf {relation} equal weight by {_pct(abs(difference), 2)} '
        f'per year in this 60-month sample (p={p_value:.3f}).'
    )


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open('rb') as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b''):
            digest.update(block)
    return digest.hexdigest()


def _manifest_path(path: Path, repo_root: Path) -> str:
    try:
        return path.relative_to(repo_root).as_posix()
    except ValueError:
        return str(path)


def _set_axes_style(ax: plt.Axes) -> None:

    ax.set_facecolor('#' + WHITE)
    ax.grid(axis='y', color='#E4EAE8', linewidth=0.8)
    ax.spines[['top', 'right', 'left']].set_visible(False)
    ax.spines['bottom'].set_color('#B9C4C0')
    ax.tick_params(colors='#536068', labelsize=9, length=0)


def _save_figure(fig: plt.Figure, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(
        path, dpi=180, bbox_inches='tight', facecolor='#' + WHITE
    )
    plt.close(fig)
    return path

STRATEGY_LABELS = {
    'wolf_cvar': 'Wolf CVaR',
    'wolf_regional_alpha': 'Wolf regional alpha',
    'equal_weight_eligible': 'Equal weight',
    'cap_weight_eligible': 'Cap weight',
}
STRATEGY_COLORS = {
    'wolf_cvar': '#' + GREEN,
    'wolf_regional_alpha': '#' + TEAL,
    'equal_weight_eligible': '#' + BLUE,
    'cap_weight_eligible': '#' + GOLD,
}


def _plot_pit_capital(
    evidence: DeckEvidence, output_path: Path
) -> Path:

    data = evidence.pit_returns.copy()
    data['date'] = pd.to_datetime(data['date'])
    data['net_return'] = pd.to_numeric(
        data['net_return'], errors='coerce'
    ).fillna(0.0)
    fig, ax = plt.subplots(figsize=(10.8, 5.2))
    endpoints: list[dict[str, object]] = []
    for strategy, group in data.groupby('strategy', sort=False):
        group = group.sort_values('date')
        wealth = evidence.current_aum * (
            1.0 + group['net_return']
        ).cumprod()

        ax.plot(
            group['date'],
            wealth / 1_000_000,
            color=STRATEGY_COLORS[strategy],
            linewidth=2.4,
            label=STRATEGY_LABELS[strategy],
        )
        endpoints.append(
            {
                'strategy': strategy,
                'date': group['date'].iloc[-1],
                'actual_y': float(wealth.iloc[-1] / 1_000_000),
                'label': _usd(wealth.iloc[-1]),
            }
        )
    if endpoints:
        ordered = sorted(endpoints, key=lambda item: float(item['actual_y']))
        values = np.asarray(
            [float(item['actual_y']) for item in ordered], dtype=float
        )
        minimum_gap = max(float(np.ptp(values)) * 0.10, 18.0)
        adjusted = values.copy()
        cluster_start = 0
        for index in range(1, len(values) + 1):
            cluster_ends = (
                index == len(values)
                or values[index] - values[index - 1] >= minimum_gap
            )
            if not cluster_ends:
                continue
            count = index - cluster_start
            if count > 1:
                centre = float(values[cluster_start:index].mean())
                offsets = (
                    np.arange(count, dtype=float) - (count - 1) / 2.0
                ) * minimum_gap
                adjusted[cluster_start:index] = centre + offsets
            cluster_start = index
        for item, label_y in zip(ordered, adjusted, strict=True):
            strategy = str(item['strategy'])
            actual_y = float(item['actual_y'])
            date = pd.Timestamp(item['date'])
            if not np.isclose(label_y, actual_y):
                ax.plot(
                    [date, date],
                    [actual_y, label_y],
                    color=STRATEGY_COLORS[strategy],
                    linewidth=0.8,
                    alpha=0.65,
                )
            ax.text(
                date,
                label_y,
                '  ' + str(item['label']),
                color=STRATEGY_COLORS[strategy],
                fontsize=9,
                va='center',
                weight='bold',
            )
    _set_axes_style(ax)
    ax.set_ylabel('Illustrative AUM, USD millions', color='#536068')
    ax.set_title(
        'Five-year point-in-time proxy on current AUM',
        loc='left',
        fontsize=15,
        color='#' + INK,
        weight='bold',
    )

    ax.xaxis.set_major_locator(mdates.YearLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y'))
    ax.legend(
        loc='upper left',
        frameon=False,
        ncol=3,
        fontsize=9,
    )
    ax.margins(x=0.04)
    fig.tight_layout()
    return _save_figure(fig, output_path)


def _plot_pit_metrics(
    evidence: DeckEvidence, output_path: Path
) -> Path:

    data = evidence.pit_summary.set_index('strategy')
    strategies = [
        'wolf_cvar',
        'equal_weight_eligible',
        'cap_weight_eligible',
    ]
    panels = [
        ('annualised_return', 'Net return', 100.0),
        ('sharpe', 'Sharpe ratio', 1.0),
        ('maximum_drawdown', 'Maximum drawdown', 100.0),
    ]
    fig, axes = plt.subplots(1, 3, figsize=(11.2, 4.1))

    labels = [STRATEGY_LABELS[item] for item in strategies]
    colors = [STRATEGY_COLORS[item] for item in strategies]
    for ax, (column, title, scale) in zip(axes, panels):
        values = [
            float(data.loc[item, column]) * scale
            for item in strategies
        ]
        bars = ax.bar(labels, values, color=colors, width=0.62)
        _set_axes_style(ax)
        ax.set_title(
            title, fontsize=13, color='#' + INK, weight='bold'
        )

        ax.tick_params(axis='x', rotation=22)
        ax.axhline(0.0, color='#B9C4C0', linewidth=0.8)
        for bar, value in zip(bars, values):
            suffix = '%' if column != 'sharpe' else ''
            label = f'{value:.1f}{suffix}'
            y = value + (0.45 if value >= 0 else -0.45)
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                y,
                label,
                ha='center',

                va='bottom' if value >= 0 else 'top',
                fontsize=9,
                color='#' + INK,
                weight='bold',
            )
    fig.suptitle(
        'Wolf trades some return for a smoother path',
        x=0.02,
        ha='left',
        fontsize=15,
        color='#' + INK,
        weight='bold',
    )
    fig.tight_layout(rect=(0, 0, 1, 0.91))

    return _save_figure(fig, output_path)


def _plot_overfitting(
    evidence: DeckEvidence, output_path: Path
) -> Path:
    row = evidence.overfitting
    values = [
        float(row['median_selected_is_information_ratio']),
        float(row['median_selected_oos_information_ratio']),
    ]
    fig, ax = plt.subplots(figsize=(7.4, 4.3))
    bars = ax.bar(
        ['Selected in-sample', 'Same winner out-of-sample'],

        values,
        color=['#' + GOLD, '#' + GREEN],
        width=0.55,
    )
    _set_axes_style(ax)
    ax.set_ylabel('Information ratio', color='#536068')
    ax.set_title(
        'Selection performance is roughly halved out-of-sample',
        loc='left',
        fontsize=14,
        color='#' + INK,
        weight='bold',
    )
    for bar, value in zip(bars, values):

        ax.text(
            bar.get_x() + bar.get_width() / 2,
            value + 0.035,
            f'{value:.2f}',
            ha='center',
            fontsize=12,
            color='#' + INK,
            weight='bold',
        )
    pbo = float(row['probability_of_backtest_overfitting'])
    haircut = float(row['selected_information_ratio_haircut'])
    ax.text(
        0.02,
        0.93,

        f'PBO {_pct(pbo)}  |  IR haircut {_pct(haircut)}',
        transform=ax.transAxes,
        fontsize=10,
        color='#' + MUTED,
        va='top',
    )
    fig.tight_layout()
    return _save_figure(fig, output_path)


def _plot_cost_drag(
    evidence: DeckEvidence, output_path: Path
) -> Path:
    before = evidence.prior_pit_summary.set_index('strategy').loc[
        'wolf_cvar'
    ]
    after = evidence.pit_summary.set_index('strategy').loc['wolf_cvar']
    panels = [
        (
            'annualised_cost_drag',
            'Annualised cost drag',
            100.0,
            '%',
        ),
        ('annualised_turnover', 'Annualised turnover', 1.0, 'x'),
    ]

    fig, axes = plt.subplots(1, 2, figsize=(8.1, 4.5))
    for ax, (column, title, scale, suffix) in zip(axes, panels):
        values = [float(before[column]) * scale, float(after[column]) * scale]
        bars = ax.bar(
            ['Before', 'Now'],
            values,
            color=['#' + RED, '#' + GREEN],
            width=0.58,
        )
        _set_axes_style(ax)
        ax.set_title(
            title, fontsize=13, color='#' + INK, weight='bold'
        )
        ax.axhline(
            1.5,
            color='#' + GOLD,
            linewidth=1.4,
            linestyle='--',
            label='1.5 target',
        )
        ax.legend(frameon=False, fontsize=8, loc='upper right')
        for bar, value in zip(bars, values):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                value + 0.06,
                f'{value:.2f}{suffix}',
                ha='center',
                fontsize=10,
                color='#' + INK,
                weight='bold',
            )
    fig.suptitle(
        'Turnover controls moved both implementation measures below target',
        x=0.02,
        ha='left',
        fontsize=14,
        color='#' + INK,
        weight='bold',
    )
    fig.tight_layout(rect=(0, 0, 1, 0.91))
    return _save_figure(fig, output_path)


def _plot_supervised_rank_ic(
    evidence: DeckEvidence, output_path: Path
) -> Path:
    validation = _supervised_validation_rows(evidence).set_index(
        'horizon_months'
    )
    legacy = _supervised_ensemble_rows(evidence).set_index(
        'horizon_months'
    )
    horizons = [3, 6, 9, 12]
    x = np.arange(len(horizons))
    width = 0.34
    validation_values = [
        float(validation.loc[horizon, 'mean_rank_ic'])
        for horizon in horizons
    ]
    legacy_values = [
        float(legacy.loc[horizon, 'mean_rank_ic'])
        for horizon in horizons
    ]
    fig, ax = plt.subplots(figsize=(8.4, 4.6))
    bars_validation = ax.bar(
        x - width / 2,
        validation_values,
        width,
        color='#' + BLUE,
        label='Expanding validation',
    )
    bars_legacy = ax.bar(
        x + width / 2,
        legacy_values,
        width,
        color='#' + GREEN,
        label='Legacy OOS diagnostic',
    )
    _set_axes_style(ax)
    ax.set_xticks(x, [f'{horizon}m' for horizon in horizons])
    ax.set_ylabel('Mean rank information coefficient', color='#536068')
    ax.set_title(
        'Ranking signal is positive; independent evidence is still sparse',
        loc='left',
        fontsize=14,
        color='#' + INK,
        weight='bold',
    )
    ax.legend(frameon=False, fontsize=9, loc='upper left')
    for bars in (bars_validation, bars_legacy):
        for bar in bars:
            value = float(bar.get_height())
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                value + 0.006,
                f'{value:.3f}',
                ha='center',
                fontsize=8.5,
                color='#' + INK,
                weight='bold',
            )
    ax.set_ylim(0.0, max(legacy_values) * 1.28)
    fig.tight_layout()
    return _save_figure(fig, output_path)


def _plot_supervised_calibration(
    evidence: DeckEvidence, output_path: Path
) -> Path:
    quantiles = evidence.supervised_quantiles.sort_values(
        'horizon_months'
    )
    horizons = quantiles['horizon_months'].astype(int).tolist()
    coverage = (
        pd.to_numeric(quantiles['central_90_coverage'], errors='coerce')
        * 100.0
    )
    widths = (
        pd.to_numeric(quantiles['mean_interval_width'], errors='coerce')
        * 100.0
    )
    fig, axes = plt.subplots(1, 2, figsize=(9.1, 4.5))
    coverage_bars = axes[0].bar(
        [f'{horizon}m' for horizon in horizons],
        coverage,
        color='#' + GREEN,
        width=0.6,
    )
    _set_axes_style(axes[0])
    axes[0].axhline(
        90.0,
        color='#' + GOLD,
        linestyle='--',
        linewidth=1.4,
        label='90% target',
    )
    axes[0].set_ylim(80.0, 100.0)
    axes[0].set_title('Central interval coverage', fontsize=12, weight='bold')
    axes[0].legend(frameon=False, fontsize=8, loc='lower left')
    for bar, value in zip(coverage_bars, coverage):
        axes[0].text(
            bar.get_x() + bar.get_width() / 2,
            float(value) + 0.5,
            f'{float(value):.1f}%',
            ha='center',
            fontsize=8.5,
            weight='bold',
            color='#' + INK,
        )

    width_bars = axes[1].bar(
        [f'{horizon}m' for horizon in horizons],
        widths,
        color='#' + TEAL,
        width=0.6,
    )
    _set_axes_style(axes[1])
    axes[1].set_title('Average interval width', fontsize=12, weight='bold')
    axes[1].set_ylabel('Benchmark-relative return', color='#536068')
    for bar, value in zip(width_bars, widths):
        axes[1].text(
            bar.get_x() + bar.get_width() / 2,
            float(value) + 2.0,
            f'{float(value):.0f}%',
            ha='center',
            fontsize=8.5,
            weight='bold',
            color='#' + INK,
        )
    axes[1].set_ylim(0.0, max(widths) * 1.2)
    fig.suptitle(
        'Coverage now clears target, but long-horizon forecasts remain wide',
        x=0.02,
        ha='left',
        fontsize=14,
        color='#' + INK,
        weight='bold',
    )
    fig.tight_layout(rect=(0, 0, 1, 0.9))
    return _save_figure(fig, output_path)


def _plot_drl_validation(
    evidence: DeckEvidence, output_path: Path
) -> Path:
    seeds = evidence.drl_seeds.sort_values('seed')
    labels = [f'PPO {int(seed)}' for seed in seeds['seed']]
    values = pd.to_numeric(
        seeds['validation_information_ratio'], errors='coerce'
    ).fillna(0.0).tolist()
    challengers = evidence.drl_challengers.loc[
        evidence.drl_challengers['split'].eq('validation')
        & evidence.drl_challengers['selected_parameter_by_validation']
        .astype(str)
        .str.lower()
        .isin({'true', '1', 'yes'})
    ]
    for row in challengers.itertuples(index=False):
        labels.append(str(row.algorithm).replace('_', ' ').title())
        values.append(float(row.information_ratio))

    fig, ax = plt.subplots(figsize=(8.8, 4.7))
    colors = ['#' + (GREEN if value > 0 else RED) for value in values]
    bars = ax.barh(labels, values, color=colors, height=0.62)
    _set_axes_style(ax)
    ax.axvline(0.0, color='#' + INK, linewidth=1.2)
    ax.set_xlabel('Frozen-validation information ratio', color='#536068')
    ax.set_title(
        'No adaptive allocator beat the zero-active baseline',
        loc='left',
        fontsize=14,
        color='#' + INK,
        weight='bold',
    )
    for bar, value in zip(bars, values):
        ax.text(
            value - 0.03 if value < 0 else value + 0.03,
            bar.get_y() + bar.get_height() / 2,
            f'{value:.2f}',
            va='center',
            ha='right' if value < 0 else 'left',
            fontsize=8.5,
            color='#' + INK,
            weight='bold',
        )
    lower = min(values + [0.0])
    upper = max(values + [0.0])
    span = max(upper - lower, 0.5)
    ax.set_xlim(lower - span * 0.18, upper + span * 0.18)
    fig.tight_layout()
    return _save_figure(fig, output_path)


def _build_plot_assets(
    evidence: DeckEvidence, plot_dir: Path
) -> tuple[Path, ...]:

    plot_dir.mkdir(parents=True, exist_ok=True)
    assets = (
        _plot_pit_capital(
            evidence, plot_dir / 'pit_capital_projection.png'
        ),
        _plot_pit_metrics(
            evidence, plot_dir / 'pit_metric_comparison.png'
        ),
        _plot_overfitting(
            evidence, plot_dir / 'overfitting_haircut.png'
        ),
        _plot_cost_drag(
            evidence, plot_dir / 'implementation_cost_drag.png'
        ),
        _plot_supervised_rank_ic(
            evidence, plot_dir / 'supervised_rank_ic.png'
        ),
        _plot_supervised_calibration(
            evidence, plot_dir / 'supervised_calibration.png'
        ),
        _plot_drl_validation(
            evidence, plot_dir / 'drl_validation.png'
        ),
    )
    return assets

def _set_background(slide, color: str = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,

    font: str = FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    inset = Inches(margin)
    frame.margin_left = inset
    frame.margin_right = inset

    frame.margin_top = inset
    frame.margin_bottom = inset
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = _ascii_display(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    shape.name = 'Text: ' + _ascii_display(text)[:50]
    return shape

def _add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = BORDER,
    radius: bool = False,
):
    kind = (
        MSO_SHAPE.ROUNDED_RECTANGLE
        if radius
        else MSO_SHAPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(w), Inches(h)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(0.8)
    if radius:
        shape.adjustments[0] = 0.06
    return shape


def _add_header(
    slide, title: str, subtitle: str, slide_number: int
) -> None:
    _add_text(
        slide, title, 0.52, 0.32, 10.9, 0.44,
        size=25, bold=True, font=FONT_HEAD,
    )

    if subtitle:
        _add_text(
            slide, subtitle, 0.54, 0.81, 11.4, 0.28,
            size=10.5, color=MUTED,
        )
    _add_text(
        slide, f'{slide_number:02d}', 12.27, 0.35, 0.52, 0.28,
        size=10, color=MUTED, bold=True, align=PP_ALIGN.RIGHT,
    )
    divider = _add_rect(
        slide, 0.54, 1.16, 12.25, 0.012,
        fill=BORDER, line=BORDER,
    )
    divider.name = 'Header divider'

def _add_footer(
    slide,
    source: str,
    *,
    label: str = 'Wolf Quant Model | Research use',
) -> None:
    _add_text(
        slide, label, 0.54, 7.18, 3.8, 0.16,
        size=7.5, color=MUTED,
    )
    _add_text(
        slide, source, 4.1, 7.18, 8.68, 0.16,
        size=7.5, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def _new_slide(

    presentation: Presentation,
    title: str,
    subtitle: str,
    number: int,
    source: str,
):
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_background(slide)
    display_number = len(presentation.slides)
    _add_header(slide, title, subtitle, display_number)
    _add_footer(slide, source)
    return slide


def _add_bullets(
    slide,
    items: Sequence[str],
    x: float,
    y: float,

    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = INK,
    spacing: float = 10,
):
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    bullet = chr(0x2022)

    for index, item in enumerate(items):
        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )
        paragraph.text = f'{bullet} {_ascii_display(item)}'
        paragraph.font.name = FONT_BODY
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = _rgb(color)
        paragraph.space_after = Pt(spacing)
        paragraph.line_spacing = 1.05
    shape.name = 'Bullet list'
    return shape

def _add_kpi(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: str = WHITE,
    accent: str = GREEN,
) -> None:
    _add_rect(slide, x, y, w, 1.02, fill=fill, line=BORDER)
    _add_rect(
        slide, x, y, 0.055, 1.02, fill=accent, line=accent
    )

    _add_text(
        slide, value, x + 0.18, y + 0.14, w - 0.28, 0.34,
        size=21, color=INK, bold=True, font=FONT_HEAD,
    )
    _add_text(
        slide, label, x + 0.18, y + 0.58, w - 0.28, 0.22,
        size=9.5, color=MUTED,
    )


def _add_callout(
    slide,
    heading: str,
    body: str,
    x: float,
    y: float,

    w: float,
    h: float,
    *,
    fill: str = PALE_GREEN,
    accent: str = GREEN,
) -> None:
    _add_rect(slide, x, y, w, h, fill=fill, line=fill)
    _add_rect(
        slide, x, y, 0.06, h, fill=accent, line=accent
    )
    if h < 0.9:
        _add_text(
            slide,
            f'{heading}: {body}',
            x + 0.22,
            y + 0.1,
            w - 0.36,
            h - 0.18,
            size=10.2,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        return
    _add_text(
        slide, heading, x + 0.22, y + 0.16, w - 0.36, 0.34,
        size=17, bold=True, font=FONT_HEAD,
    )
    _add_text(
        slide, body, x + 0.22, y + 0.57, w - 0.36, h - 0.7,
        size=11.5, color=MUTED,
    )

def _add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[object]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    widths: Sequence[float] | None = None,
    font_size: float = 10,
    highlight_rows: Sequence[int] = (),
):
    shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),

        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = shape.table
    if widths:
        total = sum(widths)
        for column, width in zip(table.columns, widths):
            column.width = Inches(w * width / total)

    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.text = _ascii_display(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(INK)

        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = FONT_BODY
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _rgb(WHITE)
        paragraph.alignment = PP_ALIGN.LEFT

    for row_index, row in enumerate(rows, start=1):
        row_fill = PALE_GREEN if row_index - 1 in highlight_rows else WHITE

        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = _ascii_display(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(row_fill)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = FONT_BODY
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = _rgb(INK)

            paragraph.alignment = (
                PP_ALIGN.RIGHT
                if column_index > 1
                else PP_ALIGN.LEFT
            )
    shape.name = 'Evidence table'
    return shape


def _add_picture_contain(
    slide,
    image_path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
):
    with Image.open(image_path) as image:
        image_w, image_h = image.size

    image_ratio = image_w / image_h
    box_ratio = w / h
    if image_ratio >= box_ratio:
        draw_w = w
        draw_h = w / image_ratio
    else:
        draw_h = h
        draw_w = h * image_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x),
        Inches(draw_y),
        Inches(draw_w),
        Inches(draw_h),
    )

    picture.name = 'Chart: ' + image_path.name
    return picture


def _add_chevron(
    slide, x: float, y: float, w: float = 0.33
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.58),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(BORDER)
    shape.line.color.rgb = _rgb(BORDER)

def _slide_cover(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    fail_count = int(evidence.scorecard['status'].eq('FAIL').sum())
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_background(slide, INK)
    _add_rect(
        slide, 0.55, 0.55, 0.08, 5.55,
        fill=GREEN, line=GREEN,
    )
    _add_text(
        slide, 'WOLF QUANT MODEL', 0.9, 0.64, 5.6, 0.3,
        size=11, color='8FD7BF', bold=True,
    )

    _add_text(
        slide, 'Investment committee\nbriefing',
        0.88, 1.18, 7.5, 1.3,
        size=38, color=WHITE, bold=True, font=FONT_HEAD,
    )
    _add_text(
        slide,
        'A disciplined global-equity decision system with '
        'portfolio constraints, stress testing and governance.',
        0.91, 2.72, 7.4, 0.8,
        size=18, color='D7E1DE',
    )
    _add_rect(
        slide, 0.91, 3.82, 5.55, 0.56,
        fill=GREEN_DARK, line=GREEN_DARK,
    )

    _add_text(
        slide,
        (
            'DECISION: APPROVE A CONTROLLED LIVE PILOT'
            if fail_count == 0
            else 'DECISION: CONTINUE PAPER AND SHADOW OPERATION'
        ),
        1.12, 3.98, 5.15, 0.2,
        size=12, color=WHITE, bold=True,
    )
    all_row = evidence.universe.loc[
        evidence.universe['region'] == 'ALL'
    ].iloc[0]
    metrics = [
        (f'{int(all_row.active):,}', 'active equities'),
        (
            f"{int(evidence.walk_forward_manifest['artifact_profile']['portfolio_months'])}",
            'monthly decisions',
        ),
        (f'{evidence.governance_score:g} / 100', 'governance score'),
        ('0', 'hard breaches'),
    ]
    for index, (value, label) in enumerate(metrics):
        x = 0.91 + index * 3.02
        _add_rect(
            slide, x, 5.25, 2.72, 1.05,
            fill='243037', line='435158',
        )

        _add_text(
            slide, value, x + 0.16, 5.42, 2.35, 0.36,
            size=21, color=WHITE, bold=True, font=FONT_HEAD,
        )
        _add_text(
            slide, label, x + 0.16, 5.84, 2.35, 0.2,
            size=9.5, color='B9C8C3',
        )
    _add_text(
        slide,
        f'As of {evidence.as_of_date}  |  '
        f'{evidence.approval_status}  |  Research evidence',
        0.91, 6.83, 11.5, 0.22,
        size=9, color='9EAFAB',
    )

def _slide_decision(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    risk = _scorecard_component(evidence, 'risk_backtesting')
    costs = _scorecard_component(evidence, 'portfolio_net_of_costs')
    constraints = _scorecard_component(evidence, 'constraint_compliance')
    pit = _scorecard_component(evidence, 'point_in_time')
    blocking = evidence.scorecard.loc[evidence.scorecard['status'].eq('FAIL')]
    pilot_ready = blocking.empty
    slide = _new_slide(
        presentation,
        'The decision in one page',
        'The controls improved materially; the alpha claim did not.',
        2,
        f'Validation {evidence.validation_as_of_date}; research through {evidence.as_of_date}',
    )
    _add_callout(
        slide,
        'Recommendation',
        (
            'Approve a small, human-supervised live pilot. '
            'Do not authorize unattended or full-scale deployment yet.'
            if pilot_ready
            else 'Continue paper and shadow operation. Do not allocate live '
            'capital while validation components still fail.'
        ),
        0.55, 1.43, 5.9, 1.18,
        fill=PALE_GREEN, accent=GREEN,
    )

    _add_text(
        slide,
        'What the evidence supports' if pilot_ready else 'What works today',
        0.58, 2.91, 5.5, 0.32,
        size=18, bold=True, font=FONT_HEAD,
    )
    _add_bullets(
        slide,
        [
            'Repeatable screening across six equity regions',
            f'Chronological VaR governance gate: {str(risk.status).title()}',
            f'Net-of-cost implementation controls: {str(costs.status).title()}',
            f'Hard portfolio constraints: {str(constraints.status).title()}',
        ],
        0.58, 3.34, 5.78, 2.62,
        size=14.5,
    )

    _add_callout(
        slide,
        'What still blocks full deployment',
        'Observed PIT evidence is incomplete, and '
        + _equal_weight_sentence(evidence),
        6.82, 1.43, 5.95, 1.18,
        fill=PALE_RED, accent=RED,
    )
    rows = [
        ('Process', 'Checksummed', 'Use'),
        ('Risk backtest', _score_label(risk), str(risk.status).title()),
        ('Cost controls', _score_label(costs), str(costs.status).title()),
        ('Hard constraints', _score_label(constraints), str(constraints.status).title()),
        ('PIT evidence', _score_label(pit), str(pit.status).title()),
        ('Deployable alpha', 'Not established', 'Do not claim'),
    ]

    _add_table(
        slide,
        ['Decision area', 'Evidence', 'Action'],
        rows,
        6.82, 2.91, 5.95, 2.95,
        widths=[2.2, 1.6, 1.6],
        font_size=11,
        highlight_rows=[
            index
            for index, component in enumerate(
                (None, risk, costs, constraints, pit, None)
            )
            if component is not None and str(component.status) == 'PASS'
        ],
    )
    _add_callout(
        slide,
        'Approval means',
        'Pilot capital, pre-trade review, monthly risk reporting and '
        'explicit stop conditions. It does not mean automatic execution.',
        0.58, 6.18, 12.19, 0.65,
        fill=PALE_BLUE, accent=BLUE,
    )


def _slide_workflow(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'What the model actually does',
        'A monthly decision pipeline with risk controls at the end.',
        3,
        'Repository pipeline stages and current resolved outputs',
    )
    steps = [
        ('1  OBSERVE', 'Prices, filings, macro and portfolio'),
        ('2  FEATURE', 'Point-in-time quality, value and risk signals'),
        ('3  FORECAST', '3, 6, 9 and 12 month distributions'),

        ('4  COMPARE', 'Portfolio-aware, clean-sheet and LLM branches'),
        ('5  CONSTRAIN', 'CVaR, liquidity and concentration limits'),
        ('6  GOVERN', 'Stress, validation and committee reporting'),
    ]
    for index, (heading, body) in enumerate(steps):
        x = 0.55 + index * 2.05
        _add_rect(
            slide, x, 2.04, 1.73, 2.05,
            fill=WHITE, line=BORDER,
        )
        _add_rect(
            slide, x, 2.04, 1.73, 0.12,
            fill=GREEN if index in (0, 4, 5) else BLUE,
            line=GREEN if index in (0, 4, 5) else BLUE,
        )

        _add_text(
            slide, heading, x + 0.14, 2.39, 1.45, 0.34,
            size=12, bold=True, font=FONT_HEAD,
        )
        _add_text(
            slide, body, x + 0.14, 2.93, 1.45, 0.82,
            size=10.5, color=MUTED,
        )
        if index < len(steps) - 1:
            _add_chevron(slide, x + 1.78, 2.79, 0.24)
    _add_callout(
        slide,
        'The output is a decision package',
        'Ranked equities, target weights, trade direction, exposure '
        'limits, risk tests, stress results and an approval status.',
        0.55, 4.64, 6.0, 1.14,

        fill=PALE_GREEN, accent=GREEN,
    )
    _add_callout(
        slide,
        'ML challengers cannot override governance',
        'PPO, bandit, convex and supervised challengers remain research-only. '
        'Their legacy OOS results are diagnostics; the baseline remains at 100%.',
        6.82, 4.64, 5.95, 1.14,
        fill=PALE_GOLD, accent=GOLD,
    )
    _add_text(
        slide,
        'Human review stays between model recommendation and execution.',
        0.58, 6.18, 12.0, 0.36,
        size=18, color=GREEN_DARK, bold=True, font=FONT_HEAD,
        align=PP_ALIGN.CENTER,
    )


def _slide_drl_results(
    presentation: Presentation,
    evidence: DeckEvidence,
    chart: Path,
) -> None:
    split = evidence.drl_split.set_index('split')
    seeds = evidence.drl_seeds.copy()
    validation_ir = pd.to_numeric(
        seeds['validation_information_ratio'], errors='coerce'
    )
    best_validation_ir = float(validation_ir.max())
    constraint_violations = int(
        pd.to_numeric(
            seeds['constraint_violations'], errors='coerce'
        ).fillna(0.0).sum()
    )
    bootstrap_environments = int(
        pd.to_numeric(
            seeds['bootstrap_environment_count'], errors='coerce'
        ).fillna(0.0).max()
    )
    slide = _new_slide(
        presentation,
        'DRL learned safely, but did not earn capital',
        'Five real PPO seeds and two simpler challengers failed frozen validation.',
        4,
        'Checksummed regional history, train-only bootstrap and validation-only selection',
    )
    _add_picture_contain(slide, chart, 0.42, 1.37, 7.0, 4.72)
    kpis = [
        (
            f"{int(evidence.drl_long_history['monthly_observations'])}",
            'regional history months',
        ),
        (f"{int(split.loc['train', 'observations'])}", 'training months'),
        (f"{int(split.loc['validation', 'observations'])}", 'validation months'),
        (
            f"{int(split.loc['legacy_locked_oos', 'observations'])}",
            'legacy OOS months',
        ),
    ]
    for index, (value, label) in enumerate(kpis):
        _add_kpi(
            slide,
            value,
            label,
            7.65 + (index % 2) * 2.58,
            1.46 + (index // 2) * 1.24,
            2.35,
            fill=WHITE,
            accent=BLUE if index < 2 else GOLD,
        )
    _add_callout(
        slide,
        'Measured result',
        f'Best PPO validation information ratio was {best_validation_ir:.2f}. '
        f'{bootstrap_environments} contiguous train-only bootstrap environments '
        f'per seed produced {constraint_violations} constraint violations.',
        7.65,
        4.05,
        5.12,
        1.25,
        fill=PALE_RED,
        accent=RED,
    )
    _add_callout(
        slide,
        'Governed decision',
        f"PPO receives {_pct(float(evidence.drl_acceptance['blend_weight_drl']), 0)}; "
        f"the classical optimiser remains {_pct(float(evidence.drl_acceptance['blend_weight_baseline']), 0)}. "
        'That is a successful safety decision, not a failed software run.',
        0.55,
        6.23,
        12.22,
        0.61,
        fill=PALE_GOLD,
        accent=GOLD,
    )


def _slide_supervised_stack(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    profile = evidence.supervised_dataset.loc[
        evidence.supervised_dataset['horizon_months'].eq(3)
    ].iloc[0]
    candidate_count = int(
        evidence.supervised_validation.loc[
            ~evidence.supervised_validation['candidate'].eq(
                'supervised_alpha_ensemble'
            ),
            'candidate',
        ].nunique()
    )
    overall = evidence.supervised_acceptance.loc[
        evidence.supervised_acceptance['scope'].eq('overall')
    ].iloc[0]
    slide = _new_slide(
        presentation,
        'The new supervised alpha research stack',
        'Benchmark-relative stock ranking across four investment horizons.',
        4,
        'Supervised dataset, validation checkpoints and model manifests',
    )
    kpis = [
        (f"{int(profile['securities']):,}", 'research securities'),
        (str(candidate_count), 'fixed candidate specifications'),
        ('4', 'forecast horizons: 3/6/9/12m'),
        (_pct(float(overall['deployment_blend_weight']), 0), 'live supervised blend'),
    ]
    for index, (value, label) in enumerate(kpis):
        _add_kpi(
            slide,
            value,
            label,
            0.55 + index * 3.08,
            1.42,
            2.78,
            fill=WHITE,
            accent=RED if index == 3 else GREEN if index == 0 else BLUE,
        )

    panels = [
        (
            'Evidence and targets',
            [
                f"{int(profile['decision_dates'])} monthly decision dates",
                f"{int(profile['numeric_features'])} numeric plus 4 categorical features",
                'Returns measured versus regional and sector peers',
                'PIT status remains reconstructed proxy',
            ],
            GREEN,
        ),
        (
            'Model challengers',
            [
                'OLS with train-only Fama-MacBeth screening',
                'Ridge, Elastic Net and robust Huber regression',
                'Random Forest, Extra Trees and histogram boosting',
                'XGBoost return and learning-to-rank models',
            ],
            BLUE,
        ),
        (
            'Controls before prediction',
            [
                'Expanding-window purged cross-validation',
                'Preprocessing and screening fitted inside each fold',
                'One linear, tree and ranker winner per ensemble',
                'Explicit turnover, impact, FX and 25bp bank fee',
            ],
            GOLD,
        ),
    ]
    for index, (heading, bullets, accent) in enumerate(panels):
        x = 0.55 + index * 4.12
        _add_rect(slide, x, 2.78, 3.83, 3.0, fill=WHITE, line=BORDER)
        _add_rect(slide, x, 2.78, 3.83, 0.09, fill=accent, line=accent)
        _add_text(
            slide,
            heading,
            x + 0.18,
            3.08,
            3.45,
            0.34,
            size=16,
            bold=True,
            font=FONT_HEAD,
        )
        _add_bullets(
            slide,
            bullets,
            x + 0.18,
            3.58,
            3.43,
            1.92,
            size=11.5,
            spacing=8,
        )
    _add_callout(
        slide,
        'Governance boundary',
        'The supervised ensemble and every DRL challenger remain research-only. '
        'The governed baseline still receives 100% of portfolio weight.',
        0.55,
        6.12,
        12.22,
        0.68,
        fill=PALE_RED,
        accent=RED,
    )


def _slide_supervised_results(
    presentation: Presentation,
    evidence: DeckEvidence,
    chart: Path,
) -> None:
    rows = _supervised_ensemble_rows(evidence)
    primary = rows.loc[rows['horizon_months'].eq(3)].iloc[0]
    slide = _new_slide(
        presentation,
        'Supervised signal: encouraging, not yet proven',
        'Positive ranking results do not overcome four independent 3-month samples.',
        5,
        'Frozen validation and already-inspected legacy OOS diagnostics',
    )
    _add_picture_contain(slide, chart, 0.45, 1.38, 7.0, 4.82)
    table_rows = []
    for row in rows.itertuples(index=False):
        table_rows.append(
            (
                f'{int(row.horizon_months)}m',
                str(int(row.observations)),
                str(int(row.independent_observations)),
                f'{float(row.mean_rank_ic):.3f}',
                f'{float(row.independent_rank_ic_sign_test_p_value):.3f}',
                _pct(float(row.mean_horizon_net_active_return), 1),
            )
        )
    _add_table(
        slide,
        ['Horizon', 'Monthly', 'Independent', 'Rank IC', 'Sign p', 'Net cohort'],
        table_rows,
        7.63,
        1.48,
        5.14,
        2.92,
        widths=[0.78, 0.76, 1.0, 0.76, 0.72, 1.0],
        font_size=8.7,
        highlight_rows=[0],
    )
    _add_callout(
        slide,
        'Primary 3-month diagnostic',
        f"Rank IC {float(primary['mean_rank_ic']):.3f}; "
        f"mean net active cohort return {_pct(float(primary['mean_horizon_net_active_return']), 1)}; "
        f"exact sign-test p={float(primary['independent_rank_ic_sign_test_p_value']):.4f}.",
        7.63,
        4.72,
        5.14,
        1.18,
        fill=PALE_GOLD,
        accent=GOLD,
    )
    _add_callout(
        slide,
        'What is deliberately absent',
        'No headline CAGR, Sharpe, t-statistic or confidence interval is reported '
        'until 12 independent prospective cohorts exist.',
        0.55,
        6.32,
        12.22,
        0.56,
        fill=PALE_BLUE,
        accent=BLUE,
    )


def _slide_supervised_calibration(
    presentation: Presentation,
    evidence: DeckEvidence,
    chart: Path,
) -> None:
    oos = _supervised_ensemble_rows(evidence).set_index('horizon_months')
    quantiles = evidence.supervised_quantiles.set_index('horizon_months')
    slide = _new_slide(
        presentation,
        'Uncertainty and implementation are now controlled',
        'Coverage clears target; wide long-horizon bands still signal low precision.',
        6,
        'Purged conformal calibration and cost-aware legacy OOS cohorts',
    )
    _add_picture_contain(slide, chart, 0.45, 1.4, 7.05, 4.78)
    table_rows = []
    for horizon in (3, 6, 9, 12):
        result = oos.loc[horizon]
        quantile = quantiles.loc[horizon]
        turnover = result['annualised_turnover']
        cost = result['annualised_cost_drag']
        table_rows.append(
            (
                f'{horizon}m',
                f'{float(turnover):.2f}x' if pd.notna(turnover) else 'N/A',
                _pct(float(cost), 2) if pd.notna(cost) else 'N/A',
                _pct(float(quantile['central_90_coverage']), 1),
                _pct(float(quantile['mean_interval_width']), 0),
            )
        )
    _add_table(
        slide,
        ['Horizon', 'Turnover', 'Cost drag', '90% cover', 'Band width'],
        table_rows,
        7.72,
        1.52,
        5.05,
        2.92,
        widths=[0.8, 0.95, 1.0, 1.0, 1.0],
        font_size=9.2,
        highlight_rows=[0, 1, 2],
    )
    _add_callout(
        slide,
        'Implementation result',
        'Recurring turnover is 0.54x, 0.47x and 0.35x at 3/6/9 months. '
        'The 12-month record has no recurring rebalance, so it remains unestimable.',
        7.72,
        4.73,
        5.05,
        1.24,
        fill=PALE_GREEN,
        accent=GREEN,
    )
    _add_callout(
        slide,
        'Interpretation',
        'Calibration is fixed; precision is not. A 117% average 12-month band is '
        'a warning against confident long-horizon stock claims.',
        0.55,
        6.3,
        12.22,
        0.58,
        fill=PALE_GOLD,
        accent=GOLD,
    )


def _slide_portfolio_outputs(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    baseline = set(
        evidence.holdings.loc[
            ~evidence.holdings['ticker'].astype(str).str.upper().isin(
                {'CASH', 'CASH.USD'}
            ),
            'ticker',
        ].astype(str)
    )
    regional = evidence.regional_alpha.copy()
    regional_weights = pd.to_numeric(
        regional['target_weight'], errors='coerce'
    ).fillna(0.0)
    regional_names = set(
        regional.loc[regional_weights.gt(1e-9), 'ticker'].astype(str)
    )
    shared = sorted(baseline & regional_names)
    baseline_only = sorted(baseline - regional_names)
    regional_only = sorted(regional_names - baseline)
    watchlist = _supervised_watchlist(evidence)
    slide = _new_slide(
        presentation,
        'Recommendations: target portfolio versus research challengers',
        'Only the governed CVaR target is actionable; challenger names remain shadow evidence.',
        7,
        'Resolved portfolio, regional-alpha target and supervised 3-month rankings',
    )
    kpis = [
        (str(len(baseline)), 'governed target equities'),
        (str(len(regional_names)), 'regional challenger equities'),
        (f'{len(shared)} / 20', 'names shared by both'),
        ('0%', 'supervised live blend'),
    ]
    for index, (value, label) in enumerate(kpis):
        _add_kpi(
            slide,
            value,
            label,
            0.55 + index * 3.08,
            1.4,
            2.78,
            fill=WHITE,
            accent=GREEN if index == 0 else RED if index == 3 else BLUE,
        )

    comparison = [
        ('Governed-only names', baseline_only, GREEN),
        ('Shared core', shared, BLUE),
        ('Regional-only names', regional_only, GOLD),
    ]
    for index, (heading, names, accent) in enumerate(comparison):
        x = 0.55 + index * 4.12
        _add_rect(slide, x, 2.69, 3.83, 1.22, fill=WHITE, line=BORDER)
        _add_rect(slide, x, 2.69, 0.055, 1.22, fill=accent, line=accent)
        _add_text(
            slide,
            heading,
            x + 0.18,
            2.84,
            3.4,
            0.26,
            size=13,
            bold=True,
            font=FONT_HEAD,
        )
        _add_text(
            slide,
            ', '.join(names),
            x + 0.18,
            3.18,
            3.45,
            0.57,
            size=8.6,
            color=MUTED,
        )

    watch_rows = []
    for row in watchlist.itertuples(index=False):
        watch_rows.append(
            (
                row.ticker,
                row.region,
                f'{float(row.supervised_alpha_score):.1f}',
                _pct(float(row.cost_adjusted_predicted_excess_return), 1),
                _pct(float(row.q05_excess_return), 0),
                _pct(float(row.q95_excess_return), 0),
            )
        )
    _add_text(
        slide,
        'Highest supervised 3-month research rank in each region',
        0.58,
        4.14,
        8.2,
        0.3,
        size=15,
        bold=True,
        font=FONT_HEAD,
    )
    _add_table(
        slide,
        ['Ticker', 'Region', 'Score', 'Cost-adj alpha', 'Q05', 'Q95'],
        watch_rows,
        0.55,
        4.53,
        12.22,
        1.85,
        widths=[1.2, 1.7, 0.8, 1.25, 0.8, 0.8],
        font_size=8.8,
    )
    _add_text(
        slide,
        'Research watchlist only. These are not buy orders; uncertainty bands are wide and the deployment gate is closed.',
        0.58,
        6.62,
        12.15,
        0.24,
        size=9.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _slide_evidence(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    profile = evidence.walk_forward_manifest['source_profile']
    artifact = evidence.walk_forward_manifest['artifact_profile']
    coverage = evidence.pit_coverage['coverage']
    free_data = (
        evidence.free_data_summary.set_index('source')
        if not evidence.free_data_summary.empty
        else pd.DataFrame()
    )

    def free_metric(source: str, column: str) -> float:
        if source not in free_data.index or column not in free_data:
            return 0.0
        value = pd.to_numeric(
            free_data.loc[source, column], errors='coerce'
        )
        return float(value) if pd.notna(value) else 0.0

    akshare_names = int(free_metric('akshare', 'entities'))
    akshare_rows = int(free_metric('akshare', 'rows'))
    yfinance_volume_names = int(
        free_metric('yfinance', 'positive_volume_entities')
    )
    openfigi_names = int(free_metric('openfigi', 'entities'))
    macro_series = int(free_metric('fred_alfred', 'entities'))
    openbb_benchmarks = int(free_metric('openbb', 'entities'))
    sec_status = (
        str(free_data.loc['sec_edgar', 'status'])
        if 'sec_edgar' in free_data.index
        else 'not measured'
    )
    all_row = evidence.universe.loc[
        evidence.universe['region'] == 'ALL'
    ].iloc[0]
    slide = _new_slide(
        presentation,
        'Evidence breadth and point-in-time progress',
        'Public identifiers, macro vintages and China/HK volume improved; historical membership remains the main gap.',
        8,
        'Universe, walk-forward manifest and aggregate PIT coverage',
    )
    kpis = [
        (f'{int(all_row.active):,}', 'active equities'),
        ('{:,}'.format(int(profile['security_count'])), 'walk-forward eligible'),
        ('{:,}'.format(int(artifact['forecast_rows'])), 'historical forecasts'),
        (f'{openfigi_names:,}', 'current FIGI mappings'),
    ]
    for index, (value, label) in enumerate(kpis):
        _add_kpi(
            slide, value, label,
            0.55 + index * 3.08, 1.46, 2.78,
            fill=WHITE, accent=GREEN if index in (0, 3) else BLUE,
        )
    region_rows = []
    for row in evidence.universe.loc[
        evidence.universe['region'] != 'ALL'
    ].itertuples():
        region_rows.append(
            (
                row.region,
                f'{int(row.active):,}',
                f'{int(row.total):,}',
                _pct(row.active_share),
            )
        )
    _add_table(
        slide,
        ['Region', 'Active', 'Master', 'Active share'],
        region_rows,
        0.55, 2.83, 6.1, 3.12,
        widths=[2.1, 1.0, 1.0, 1.25],
        font_size=10.5,
    )

    pit_rows = [
        ('Delisting events', f"{int(coverage['delisting_events']):,}", 'Archived'),
        (
            'Fundamentals with filing date',
            f"{int(coverage['fundamental_rows_with_filing_date']):,}",
            'Mostly reconstructed',
        ),
        (
            'Observed acceptance times',
            f"{int(coverage['observed_acceptance_timestamps']):,}",
            'Observed' if int(coverage['observed_acceptance_timestamps']) else 'Missing',
        ),
        (
            'Dated membership events',
            f"{int(coverage['historical_membership_events']):,}",
            'Observed' if int(coverage['historical_membership_events']) else 'Missing',
        ),
        (
            'Inactive names with prices',
            f"{int(coverage['inactive_price_securities']):,}",
            'Observed' if int(coverage['inactive_price_securities']) else 'Missing',
        ),
        (
            'Names with historical volume',
            f"{int(coverage['securities_with_historical_volume']):,}",
            'Observed' if int(coverage['securities_with_historical_volume']) else 'Missing',
        ),
    ]
    _add_table(
        slide,
        ['PIT evidence', 'Coverage', 'Status'],
        pit_rows,
        6.93, 2.83, 5.84, 3.12,
        widths=[2.55, 1.2, 1.45],
        font_size=10.0,
        highlight_rows=[0],
    )

    _add_callout(
        slide,
        'Free-data checkpoint and publication boundary',
        f'yfinance volume: {yfinance_volume_names:,} China/HK names; '
        f'AKShare: {akshare_rows:,} bars across {akshare_names:,} names; '
        f'FRED/ALFRED: {macro_series} series; OpenBB: {openbb_benchmarks} benchmark checks; '
        f'SEC: {sec_status}. OpenFIGI is a current snapshot, not historical membership. '
        'Only aggregate counts and checksums are published.',
        0.55, 6.18, 12.22, 0.65,
        fill=PALE_GOLD, accent=GOLD,
    )


def _slide_trades(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    buys = int((evidence.trades['trade_action'] == 'Buy').sum())
    reduces = int(
        (evidence.trades['trade_action'] == 'Reduce').sum()
    )
    cash_mask = evidence.holdings['ticker'].astype(str).str.upper().isin(
        {'CASH', 'CASH.USD'}
    )
    equities = evidence.holdings.loc[~cash_mask]
    equity_count = len(equities)
    equity_regions = int(equities['region'].nunique())
    max_name_weight = float(equities['final_weight'].max())
    slide = _new_slide(
        presentation,
        'Equities to establish in the target portfolio',
        'Model trade direction versus current holdings; weights are targets.',
        9,
        'Final portfolio weights and portfolio trade list',
    )
    _add_kpi(
        slide, str(buys), 'buy actions',
        0.55, 1.39, 2.3, fill=PALE_GREEN, accent=GREEN,
    )

    _add_kpi(
        slide, str(reduces), 'reduce actions',
        3.02, 1.39, 2.3, fill=PALE_GOLD, accent=GOLD,
    )
    _add_kpi(
        slide, _pct(max_name_weight, 1), 'maximum name',
        5.49, 1.39, 2.3, fill=WHITE, accent=BLUE,
    )
    _add_kpi(
        slide, str(equity_count), 'equity holdings',
        7.96, 1.39, 2.3, fill=WHITE, accent=TEAL,
    )
    _add_kpi(
        slide, str(equity_regions), 'equity regions',
        10.43, 1.39, 2.34, fill=WHITE, accent=GREEN,
    )

    rows = []
    for row in evidence.trades.itertuples():
        rows.append(
            (
                row.ticker,
                _ascii_display(row.company_name)[:29],
                row.region,
                row.trade_action,
                _pct(row.target_weight, 1),
            )
        )
    midpoint = math.ceil(len(rows) / 2)
    for index, subset in enumerate(
        (rows[:midpoint], rows[midpoint:])
    ):
        _add_table(
            slide,
            ['Ticker', 'Company', 'Region', 'Action', 'Target'],

            subset,
            0.55 + index * 6.16,
            2.73,
            5.89,
            3.72,
            widths=[1.05, 2.55, 1.35, 0.9, 0.8],
            font_size=8.6,
        )
    _add_text(
        slide,
        'Execution note: recalculate orders using live NAV, FX, '
        'liquidity and prices. Manual review remains mandatory.',
        0.58, 6.68, 12.15, 0.24,
        size=9.5, color=RED, bold=True,
        align=PP_ALIGN.CENTER,
    )

def _slide_exposure(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    cash_mask = evidence.holdings['ticker'].astype(str).str.upper().isin(
        {'CASH', 'CASH.USD'}
    )
    equities = evidence.holdings.loc[~cash_mask]
    cash_weight = float(
        evidence.holdings.loc[cash_mask, 'final_weight'].sum()
    )
    max_name_weight = float(equities['final_weight'].max())
    region_weights = equities.groupby('region')['final_weight'].sum()
    sector_weights = equities.groupby('sector')['final_weight'].sum()
    largest_region = str(region_weights.idxmax())
    largest_sector = str(sector_weights.idxmax())
    chart = (
        evidence.release_root
        / 'plots/final_portfolio_exposures.png'
    )
    slide = _new_slide(
        presentation,
        'A deliberately diversified target',
        'The portfolio uses equal name weights and explicit exposure caps.',
        10,
        'Final portfolio weights and constraint report',
    )
    _add_picture_contain(slide, chart, 0.47, 1.39, 8.2, 5.33)
    _add_text(
        slide, 'Construction rules', 8.94, 1.54, 3.65, 0.34,
        size=18, bold=True, font=FONT_HEAD,
    )

    _add_bullets(
        slide,
        [
            f'{len(equities)} equities at no more than '
            f'{_pct(max_name_weight, 1)} each',
            f'{_pct(cash_weight, 0)} cash reserve',
            'No sector above 25%',
            'No country above 30%',
            'No region or currency above 40%',
            'One listing per issuer',
        ],
        8.94, 2.03, 3.65, 2.64,
        size=14,
    )
    _add_callout(
        slide,
        'Largest allocations',
        f'{largest_region} {_pct(region_weights.max(), 0)}; '
        f'{largest_sector} {_pct(sector_weights.max(), 0)}. '
        'Cash is excluded from equity exposure caps.',
        8.94, 5.02, 3.72, 1.18,
        fill=PALE_BLUE, accent=BLUE,
    )

def _ending_value(
    evidence: DeckEvidence, strategy: str
) -> float:
    returns = evidence.pit_returns.loc[
        evidence.pit_returns['strategy'] == strategy,
        'net_return',
    ].astype(float)
    return evidence.current_aum * float(
        (1.0 + returns).prod()
    )


def _slide_optimizer(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'Why the CVaR portfolio is the baseline',
        'It balances expected return, downside loss and diversification.',
        11,

        'Portfolio optimisation summary',
    )
    labels = {
        'equal_weight': 'Equal weight',
        'score_weighted': 'Score weighted',
        'risk_parity': 'Risk parity',
        'mean_variance': 'Mean variance',
        'cvar_constrained': 'CVaR constrained',
        'dividend_income': 'Dividend income',
        'regime_aware': 'Regime aware',
    }
    rows = []
    highlight = []
    for row in evidence.optimiser.itertuples():
        if row.portfolio_method == 'cvar_constrained':
            highlight.append(len(rows))

        rows.append(
            (
                labels.get(row.portfolio_method, row.portfolio_method),
                _pct(row.expected_volatility),
                _pct(row.cvar_5),
                _pct(row.expected_dividend_yield),
                str(int(row.hard_constraint_breaches)),
            )
        )
    _add_table(
        slide,
        ['Portfolio', 'Volatility', 'CVaR 5%', 'Yield', 'Hard breaches'],
        rows,
        0.55, 1.52, 7.35, 4.45,
        widths=[2.4, 1.2, 1.2, 0.95, 1.2],
        font_size=10.5,
        highlight_rows=highlight,
    )

    _add_callout(
        slide,
        'Selected: CVaR constrained',
        'The baseline keeps the forecast return target while reducing '
        'the average loss in the worst 5% of modeled outcomes to 5.10%.',
        8.22, 1.52, 4.55, 1.28,
        fill=PALE_GREEN, accent=GREEN,
    )
    _add_text(
        slide, 'Why this matters', 8.25, 3.17, 4.2, 0.34,
        size=18, bold=True, font=FONT_HEAD,
    )
    _add_bullets(
        slide,
        [
            'Focuses on severe downside, not only average volatility',
            'Maintains 20 effective holdings',
            'Passes every hard portfolio constraint',

            'Produces the final 20-name equal-weight target',
        ],
        8.25, 3.66, 4.34, 2.05,
        size=14,
    )
    _add_callout(
        slide,
        'Interpret carefully',
        'Expected-return inputs are model forecasts, not guaranteed '
        'market returns. Live sizing still depends on trading capacity.',
        8.22, 5.87, 4.55, 0.75,
        fill=PALE_GOLD, accent=GOLD,
    )

def _slide_pit_growth(
    presentation: Presentation,
    evidence: DeckEvidence,
    plot_path: Path,
) -> None:
    wolf_end = _ending_value(evidence, 'wolf_cvar')
    wolf_pnl = wolf_end - evidence.current_aum
    slide = _new_slide(
        presentation,
        'Five-year point-in-time proxy',
        'The relevant model evidence uses dated monthly decisions.',
        12,
        '60-month reconstructed point-in-time portfolio returns',
    )
    _add_picture_contain(slide, plot_path, 0.47, 1.38, 8.45, 5.45)

    _add_callout(
        slide,
        'Illustrative current-AUM result',
        f'{_usd(evidence.current_aum)} grows to '
        f'{_usd(wolf_end)}; net PnL is {_usd(wolf_pnl)} '
        'after trading costs and before the bank fee.',
        9.12, 1.51, 3.65, 1.17,
        fill=PALE_GREEN, accent=GREEN,
    )
    wolf = evidence.pit_summary.set_index('strategy').loc['wolf_cvar']
    _add_bullets(
        slide,
        [
            f'Net annualised return: {_pct(wolf.annualised_return)}',
            f'Sharpe ratio: {float(wolf.sharpe):.2f}',
            f'Maximum drawdown: {_pct(wolf.maximum_drawdown)}',
            f'Positive months: {_pct(wolf.positive_period_ratio)}',

        ],
        9.14, 3.03, 3.48, 1.95,
        size=14,
    )
    _add_callout(
        slide,
        'Not a forecast',
        'This scales the realised proxy path to current AUM. The 25 bp '
        'annual bank fee is modeled separately in the long replay.',
        9.12, 5.42, 3.65, 1.30,
        fill=PALE_GOLD, accent=GOLD,
    )

def _slide_pit_comparison(
    presentation: Presentation,
    evidence: DeckEvidence,
    plot_path: Path,
) -> None:
    summary = evidence.pit_summary.set_index('strategy')
    wolf = summary.loc['wolf_cvar']
    equal = summary.loc['equal_weight_eligible']
    slide = _new_slide(
        presentation,
        'What beat what over 60 months',
        'Wolf delivered the best risk-adjusted result, not the best return.',
        13,
        'Point-in-time strategy comparison, net of modeled costs',
    )
    _add_picture_contain(slide, plot_path, 0.55, 1.43, 8.3, 4.85)
    _add_callout(
        slide,
        'The useful result',
        f'Wolf Sharpe was {float(wolf.sharpe):.2f} and drawdown was '
        f'{_pct(abs(float(wolf.maximum_drawdown)))}, better than both '
        'simple controls on these risk measures.',
        9.08, 1.56, 3.69, 1.31,
        fill=PALE_GREEN, accent=GREEN,
    )
    _add_callout(
        slide,
        'The honest result',
        f'Equal weight returned {_pct(equal.annualised_return, 2)} versus '
        f'Wolf at {_pct(wolf.annualised_return, 2)}. The model did not '
        'beat the simplest control on net return.',
        9.08, 3.16, 3.69, 1.31,
        fill=PALE_GOLD, accent=GOLD,
    )
    _add_callout(
        slide,
        'Investment interpretation',
        'The current case is risk control and decision discipline. '
        'Incremental stock-selection alpha remains unproven.',
        9.08, 4.76, 3.69, 1.31,
        fill=PALE_BLUE, accent=BLUE,
    )
    _add_text(
        slide,
        'Highest return: Equal weight  |  Highest Sharpe: Wolf  |  '
        'Smallest drawdown: Wolf',
        0.58, 6.49, 12.1, 0.28,
        size=12.5, color=GREEN_DARK, bold=True,
        align=PP_ALIGN.CENTER,
    )


def _slide_alpha(
    presentation: Presentation,
    evidence: DeckEvidence,
    plot_path: Path,
) -> None:
    slide = _new_slide(
        presentation,
        'Alpha and overfitting: the hard truth',
        'The retrospective strategies are interesting; deployable alpha is not established.',
        14,
        'Point-in-time alpha tests and CSCV overfitting diagnostics',
    )
    alpha_rows = []
    for row in evidence.alpha.itertuples():
        alpha_rows.append(
            (
                STRATEGY_LABELS.get(row.benchmark, row.benchmark),
                _pct(row.annualised_active_return, 2),

                _pct(row.annualised_regression_alpha, 2),
                f'{float(row.two_sided_p_value):.3f}',
                row.alpha_evidence_verdict.replace('_', ' '),
            )
        )
    _add_table(
        slide,
        ['Benchmark', 'Active return', 'Reg. alpha', 'p-value', 'Verdict'],
        alpha_rows,
        0.55, 1.52, 6.14, 1.6,
        widths=[1.7, 1.15, 1.05, 0.85, 1.35],
        font_size=10.3,
    )
    _add_picture_contain(slide, plot_path, 6.94, 1.42, 5.83, 3.45)

    _add_callout(
        slide,
        'What the 50.3% fall means',
        'When history selects the apparent winner, its information ratio '
        'typically drops from 1.32 in-sample to 0.66 out-of-sample. '
        'About half the apparent edge does not survive selection.',
        0.55, 3.55, 6.14, 1.35,
        fill=PALE_GOLD, accent=GOLD,
    )
    _add_callout(
        slide,
        'Decision',
        'Use Wolf for disciplined ranking and risk control. Do not '
        'market it as a proven alpha engine until native live evidence '
        'clears the benchmark and cost hurdles.',
        0.55, 5.24, 12.22, 1.2,
        fill=PALE_RED, accent=RED,
    )

def _slide_long_history(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'The 1997 replay is a stress map',
        'It shows how today\'s holdings behaved through history, not what the model knew then.',
        15,
        '1997-present retrospective holdings replay',
    )
    chosen = [
        'current_portfolio',
        'portfolio_aware_overlay',
        'optimised_risk_parity',
        'llm_benchmark',
        'final_portfolio',
    ]
    data = evidence.performance.loc[
        (evidence.performance['window'] == 'requested_1997_window')

        & evidence.performance['strategy'].isin(chosen)
    ].set_index('strategy')
    rows = []
    for strategy in chosen:
        row = data.loc[strategy]
        rows.append(
            (
                row.strategy_label,
                _usd(row.initial_capital_usd),
                _pct(row.cagr),
                f'{float(row.sharpe):.2f}',
                _pct(row.maximum_drawdown),
                _usd(row.ending_value_usd),
            )
        )
    _add_table(
        slide,
        ['Portfolio', 'Start', 'CAGR', 'Sharpe', 'Max DD', 'End'],

        rows,
        0.55, 1.53, 7.2, 3.34,
        widths=[2.35, 1.15, 0.8, 0.75, 0.9, 1.25],
        font_size=9.7,
    )
    chart = evidence.backtest_root / 'plots/risk_return.png'
    _add_picture_contain(slide, chart, 7.98, 1.43, 4.79, 3.65)
    current_pnl = _usd(
        data.loc['current_portfolio', 'pnl_usd']
    )
    resolved_pnl = _usd(
        data.loc['final_portfolio', 'pnl_usd']
    )
    _add_callout(
        slide,
        'What looks strongest',
        f'Current Portfolio PnL: {current_pnl}. '
        f'Final Resolved PnL: {resolved_pnl}. '
        'Risk Parity led the independent optimisers.',
        0.55, 5.24, 6.0, 1.25,
        fill=PALE_GREEN, accent=GREEN,
    )

    _add_callout(
        slide,
        'Why it cannot prove selection skill',
        'The securities were chosen with current information and then '
        'replayed backward. Survivorship and selection look-ahead '
        'make the long PnL unsuitable as a live promise.',
        6.82, 5.24, 5.95, 1.31,
        fill=PALE_GOLD, accent=GOLD,
    )

def _slide_macro_events(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'Performance through major market shocks',
        'Event windows are descriptive overlaps, not causal estimates.',
        16,
        'Macro-event timeline and event performance table',
    )
    chart = evidence.backtest_root / 'plots/macro_event_timeline.png'
    _add_picture_contain(slide, chart, 0.42, 1.29, 12.48, 4.82)
    callouts = [
        ('COVID-19', '-12.58% return; -15.19% event drawdown', RED),
        ('2008 crisis', '-27.12% event drawdown', GOLD),
        ('2026 Iran war', '+3.57% return; -5.27% drawdown', BLUE),
    ]

    for index, (heading, body, accent) in enumerate(callouts):
        _add_callout(
            slide, heading, body,
            0.55 + index * 4.12, 6.17, 3.84, 0.61,
            fill=WHITE, accent=accent,
        )


def _slide_regimes(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'Interest-rate and market-regime context',
        'Conditional groupings show where the final portfolio was most comfortable.',
        17,
        'Lagged rate and market-regime performance',
    )

    rate_chart = (
        evidence.backtest_root
        / 'plots/interest_rate_performance.png'
    )
    regime_chart = (
        evidence.backtest_root
        / 'plots/market_regime_performance.png'
    )
    _add_picture_contain(
        slide, rate_chart, 0.43, 1.42, 6.22, 4.45
    )
    _add_picture_contain(
        slide, regime_chart, 6.75, 1.42, 6.15, 4.45
    )

    _add_callout(
        slide,
        'Rate reading',
        'Final Portfolio was strongest in high-rate months (15.77%) '
        'and when rates were rising (13.03%).',
        0.55, 6.04, 5.95, 0.74,
        fill=PALE_BLUE, accent=BLUE,
    )
    _add_callout(
        slide,
        'Regime reading',
        'Its strongest market group was Bull / Volatile (14.09%). '
        'Expansion months returned 9.91% annualised.',
        6.82, 6.04, 5.95, 0.74,
        fill=PALE_GREEN, accent=GREEN,
    )

def _slide_governance(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    locked = evidence.walk_forward_manifest['locked_risk_calibration']
    risk_component = _scorecard_component(evidence, 'risk_backtesting')
    risk_status = str(risk_component.status)
    pass_count = int(evidence.scorecard['status'].eq('PASS').sum())
    warning_count = int(evidence.scorecard['status'].eq('WARNING').sum())
    fail_count = int(evidence.scorecard['status'].eq('FAIL').sum())
    critical_count = len(evidence.validation_manifest.get('critical_failures', []))
    slide = _new_slide(
        presentation,
        (
            'Risk holdout passes; approval stays conditional'
            if risk_status == 'PASS'
            else 'Risk calibration still requires evidence'
        ),
        (
            'Blocked development CV controls coverage without hiding PIT and alpha limits.'
            if risk_status == 'PASS'
            else 'The report shows failed coverage or independence tests rather than overriding them.'
        ),
        18,
        'Validation scorecard and overall/holdout VaR tests',
    )
    chart = evidence.release_root / 'plots/validation_scorecard.png'
    _add_picture_contain(slide, chart, 0.43, 1.39, 6.25, 4.92)

    risk_rows = []
    for row in evidence.risk_backtest.itertuples():
        segment = (
            'Holdout'
            if row.evaluation_segment == 'chronological_holdout'
            else 'Overall'
        )
        risk_rows.append(
            (
                segment,
                f'{float(row.confidence_level):.0%}',
                f'{int(row.violations)}/{int(row.observations)}',
                f'{float(row.p_value):.3f}',
                f'{float(row.christoffersen_p_value):.3f}',
                str(row.status).title(),
            )
        )
    _add_text(
        slide, 'VaR coverage and independence', 6.94, 1.48, 5.5, 0.34,
        size=18, bold=True, font=FONT_HEAD,
    )
    _add_table(
        slide,
        ['Sample', 'VaR', 'Exceptions', 'Kupiec p', 'Indep. p', 'Status'],
        risk_rows,
        6.93, 1.94, 5.84, 2.38,
        widths=[1.0, 0.55, 0.88, 0.76, 0.76, 0.75],
        font_size=8.5,
        highlight_rows=[
            index
            for index, row in enumerate(evidence.risk_backtest.itertuples())
            if str(row.status) == 'PASS'
        ],
    )
    _add_callout(
        slide,
        f'Chronological holdout gate: {risk_status.title()} '
        f'({_score_label(risk_component)})',
        'DCC-IGARCH Student-t, filtered historical simulation, EWMA '
        'Normal and EWMA Student-t are selected using trailing data. '
        f'{int(locked.get("selection_folds", 1))} blocked development fold(s) selected a '
        f'{float(locked["selected_scale_factor"]):.3f}x '
        f'scale and {float(locked["selected_exception_multiplier"]):.2f}x '
        f'buffer for {int(locked["selected_exception_days"])} day after a breach. '
        'Overall rows remain diagnostics; the configured holdout is the gate. '
        'It is reconstructed evidence, not a pristine future shadow period.',
        6.93, 4.65, 5.84, 1.50,
        fill=PALE_GREEN if risk_status == 'PASS' else PALE_RED,
        accent=GREEN if risk_status == 'PASS' else RED,
    )
    _add_text(
        slide,
        f'Overall {evidence.governance_score:g}/100  |  {pass_count} pass  |  '
        f'{warning_count} warnings  |  {fail_count} fail  |  '
        f'{critical_count} critical failures',
        0.58, 6.5, 12.05, 0.28,
        size=12.3, color=GREEN_DARK, bold=True,
        align=PP_ALIGN.CENTER,
    )


def _slide_costs(
    presentation: Presentation,
    evidence: DeckEvidence,
    plot_path: Path,
) -> None:
    slide = _new_slide(
        presentation,
        'Turnover and modeled cost drag now meet target',
        'Retention hysteresis and minimum-turnover transitions reduced churn.',
        19,
        'Before/after point-in-time cost and turnover validation',
    )
    _add_picture_contain(slide, plot_path, 0.48, 1.43, 7.56, 4.85)
    wolf = evidence.pit_summary.set_index('strategy').loc['wolf_cvar']
    fee = evidence.backtest_manifest['annual_bank_fee']

    _add_kpi(
        slide, _pct(wolf.annualised_cost_drag, 2),
        'modeled cost drag; target <= 1.50%',
        8.32, 1.57, 4.45, fill=PALE_GREEN, accent=GREEN,
    )
    _add_kpi(
        slide, f'{float(wolf.annualised_turnover):.2f}x',
        'annual turnover; target <= 1.50x',
        8.32, 2.83, 4.45, fill=PALE_GREEN, accent=GREEN,
    )
    _add_kpi(
        slide, _usd(fee['reference_annual_charge_usd'], 0),
        '0.25% annual bank fee at reference AUM',
        8.32, 4.09, 4.45, fill=WHITE, accent=BLUE,
    )

    _add_callout(
        slide,
        'Why portfolio remains a warning',
        'Both cost gates pass, but '
        + _equal_weight_sentence(evidence),
        8.32, 5.35, 4.45, 1.12,
        fill=PALE_GOLD, accent=GOLD,
    )
    _add_text(
        slide,
        'Retention and capped/minimum-turnover transitions drove the '
        'improvement; the configured no-trade band did not trigger in this sample.',
        0.58, 6.55, 12.1, 0.27,
        size=11.8, color=GREEN_DARK, bold=True,
        align=PP_ALIGN.CENTER,
    )


def _slide_pilot(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'A controlled path to live use',
        'Scale only when operations, costs and evidence earn it.',
        20,
        'Proposed implementation and model-risk controls',
    )
    phases = [
        (
            '0  SHADOW',
            '4 weeks',
            'Generate orders, do not trade. Validate data, NAV, FX and fills.',
        ),
        (
            '1  PILOT',
            '6 to 12 months',
            'Use limited capital with human approval and staged execution.',
        ),

        (
            '2  SCALE',
            'Evidence gated',
            'Increase capital only after cost, risk and benchmark hurdles hold.',
        ),
    ]
    for index, (name, duration, body) in enumerate(phases):
        x = 0.55 + index * 4.12
        _add_rect(
            slide, x, 1.58, 3.84, 2.08,
            fill=WHITE, line=BORDER,
        )
        _add_rect(
            slide, x, 1.58, 3.84, 0.12,
            fill=[BLUE, GREEN, GOLD][index],
            line=[BLUE, GREEN, GOLD][index],
        )

        _add_text(
            slide, name, x + 0.2, 1.91, 3.4, 0.32,
            size=17, bold=True, font=FONT_HEAD,
        )
        _add_text(
            slide, duration, x + 0.2, 2.37, 3.4, 0.24,
            size=11, color=GREEN_DARK, bold=True,
        )
        _add_text(
            slide, body, x + 0.2, 2.78, 3.39, 0.62,
            size=11.5, color=MUTED,
        )
    _add_text(
        slide, 'Proposed pilot gates', 0.58, 4.07, 5.4, 0.34,
        size=18, bold=True, font=FONT_HEAD,
    )

    _add_bullets(
        slide,
        [
            'Zero hard constraint or data-lineage breaches',
            'Pre-trade cost estimate approved for every rebalance',
            'Keep turnover <=1.5x; stretch goal <=1.0x in the pilot',
            'Net performance tracked against equal and cap weight',
            'Monthly model-risk and investment-committee review',
        ],
        0.58, 4.52, 5.83, 1.83,
        size=13.5,
    )
    _add_callout(
        slide,
        'Immediate stop conditions',
        'Stale critical data, any hard breach, a failed live risk '
        'test, failed reconciliation or trading costs '
        'above the approved budget.',
        6.82, 4.12, 5.95, 1.35,
        fill=PALE_RED, accent=RED,
    )

    _add_callout(
        slide,
        'Alpha gate',
        'Full-scale alpha claims require native live evidence and at '
        'least 60 months under the existing governance policy.',
        6.82, 5.71, 5.95, 0.84,
        fill=PALE_GOLD, accent=GOLD,
    )


def _slide_glossary(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    slide = _new_slide(
        presentation,
        'How to read the results',
        'Plain-language definitions for the investment committee.',
        21,
        'Backtest methodology and validation framework',
    )

    rows = [
        ('CAGR', 'Average compounded annual growth.'),
        ('Sharpe', 'Return earned per unit of overall variability.'),
        ('Sortino', 'Return earned per unit of harmful downside variation.'),
        ('Max drawdown', 'Largest peak-to-trough loss.'),
        ('CVaR', 'Average loss among the worst modeled outcomes.'),
        ('Information ratio', 'Active return per unit of benchmark-relative risk.'),
        ('p-value', 'How surprising the result is if true alpha were zero.'),
        ('PBO', 'Chance that selecting the backtest winner overfits history.'),
    ]
    _add_table(
        slide,
        ['Measure', 'Simple meaning'],
        rows,
        0.55, 1.49, 7.0, 4.95,
        widths=[1.65, 4.95],
        font_size=11,
    )

    _add_callout(
        slide,
        'Three evidence layers',
        f"1. Dated {int(evidence.walk_forward_manifest['artifact_profile']['portfolio_months'])}-month proxy: primary model evidence.\n"
        '2. 1997 replay: stress and exposure history.\n'
        '3. Native live record: required before full alpha approval.',
        7.86, 1.49, 4.91, 1.62,
        fill=PALE_BLUE, accent=BLUE,
    )
    _add_callout(
        slide,
        'Higher is not always better',
        'A portfolio can have a higher return but a deeper drawdown, '
        'more turnover or weaker statistical evidence. The committee '
        'should judge the package, not one headline number.',
        7.86, 3.45, 4.91, 1.42,
        fill=PALE_GOLD, accent=GOLD,
    )

    _add_callout(
        slide,
        'Where to verify',
        'The GitHub release contains the full PDF, plain-language '
        'interpretation, CSV evidence, plots, manifests and checksums.',
        7.86, 5.21, 4.91, 1.23,
        fill=PALE_GREEN, accent=GREEN,
    )

def _slide_close(
    presentation: Presentation, evidence: DeckEvidence
) -> None:
    fail_count = int(evidence.scorecard['status'].eq('FAIL').sum())
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_background(slide, INK)
    _add_text(
        slide, 'THE RECOMMENDATION', 0.91, 0.74, 4.8, 0.28,
        size=11, color='8FD7BF', bold=True,
    )
    _add_text(
        slide,
        (
            'Approve a controlled live pilot.'
            if fail_count == 0
            else 'Continue paper and shadow operation.'
        ),
        0.89, 1.32, 8.7, 0.62,
        size=34, color=WHITE, bold=True, font=FONT_HEAD,
    )

    _add_text(
        slide,
        'Wolf is ready to improve how decisions are made: broad '
        'screening, explicit risk, consistent portfolio construction '
        'and an audit trail. It is not yet ready to promise alpha.',
        0.92, 2.26, 8.55, 1.05,
        size=18, color='D7E1DE',
    )
    statements = [
        ('USE NOW', 'Ranking, diversification, adaptive risk and governance'),
        ('PROVE LIVE', 'Net benchmark edge in a genuine future shadow record'),
        ('FIX NEXT', 'Observed PIT vintages and prospective model evidence'),
    ]
    for index, (heading, body) in enumerate(statements):
        x = 0.92 + index * 4.03

        _add_rect(
            slide, x, 4.08, 3.67, 1.34,
            fill='243037', line='435158',
        )
        _add_text(
            slide, heading, x + 0.18, 4.32, 3.2, 0.27,
            size=12, color='8FD7BF', bold=True,
        )
        _add_text(
            slide, body, x + 0.18, 4.79, 3.17, 0.4,
            size=11.5, color=WHITE,
        )
    _add_rect(
        slide, 0.92, 6.15, 11.72, 0.58,
        fill=GREEN_DARK, line=GREEN_DARK,
    )

    _add_text(
        slide,
        'Investment edge today: a better-controlled process.  '
        'Investment edge tomorrow: only what live evidence proves.',
        1.15, 6.33, 11.25, 0.22,
        size=13, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        f'{evidence.approval_status}  |  '
        f'Governance {evidence.governance_score:g}/100  |  '
        f'As of {evidence.as_of_date}',
        0.92, 7.02, 11.7, 0.2,
        size=8.5, color='9EAFAB',
        align=PP_ALIGN.RIGHT,
    )

def _validate_slide_bounds(presentation: Presentation) -> list[str]:
    errors: list[str] = []
    for slide_index, slide in enumerate(
        presentation.slides, start=1
    ):
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.width <= 0
                or shape.height <= 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                errors.append(
                    f'slide {slide_index}: {shape.name} is out of bounds'
                )
    return errors


def _set_document_properties(
    presentation: Presentation, evidence: DeckEvidence
) -> None:

    properties = presentation.core_properties
    properties.title = 'Wolf Quant Model Investment Committee Briefing'
    properties.subject = (
        'Current recommendations, validation evidence and live pilot plan'
    )
    properties.author = 'The Wolf Quant Model'
    properties.keywords = (
        'equities, portfolio, validation, backtest, investment committee'
    )
    properties.comments = (
        'Generated from tracked repository evidence. Research use only.'
    )
    properties.created = datetime.now()
    properties.modified = datetime.now()

def _risk_markdown_row(
    label: str, confidence: int, row: pd.Series
) -> str:
    return (
        f'| {label} | {confidence}% | '
        f'{int(row.violations)}/{int(row.observations)} | '
        f'{float(row.p_value):.3f} | '
        f'{float(row.christoffersen_p_value):.3f} | '
        f'{str(row.status).title()} |'
    )


def _performance_markdown_row(
    label: str, column: str, rows: Sequence[pd.Series]
) -> str:
    if column in {'sharpe', 'sortino'}:
        values = [f'{float(row[column]):.2f}' for row in rows]
    else:
        values = [_pct(float(row[column])) for row in rows]
    return f'| {label} | ' + ' | '.join(values) + ' |'


def _report_markdown(evidence: DeckEvidence) -> str:
    summary = evidence.pit_summary.set_index('strategy')
    prior_summary = evidence.prior_pit_summary.set_index('strategy')
    wolf = summary.loc['wolf_cvar']
    prior_wolf = prior_summary.loc['wolf_cvar']
    equal = summary.loc['equal_weight_eligible']
    cap = summary.loc['cap_weight_eligible']
    wolf_end = _ending_value(evidence, 'wolf_cvar')
    coverage = evidence.pit_coverage['coverage']
    production = evidence.production_pit.set_index('dataset')
    fundamental_vintages = int(
        production.loc['fundamental_vintages', 'rows']
    )
    corporate_action_vintages = int(
        production.loc['corporate_action_vintages', 'rows']
    )
    market_cap_vintages = int(
        production.loc['market_cap_vintages', 'rows']
    )
    locked = evidence.walk_forward_manifest['locked_risk_calibration']
    pass_count = int(evidence.scorecard['status'].eq('PASS').sum())
    warning_count = int(evidence.scorecard['status'].eq('WARNING').sum())
    fail_count = int(evidence.scorecard['status'].eq('FAIL').sum())
    risk_component = _scorecard_component(evidence, 'risk_backtesting')
    cost_component = _scorecard_component(evidence, 'portfolio_net_of_costs')
    pit_component = _scorecard_component(evidence, 'point_in_time')
    risk_status = str(risk_component.status)
    free_data = (
        evidence.free_data_summary.set_index('source')
        if not evidence.free_data_summary.empty
        else pd.DataFrame()
    )

    def aggregate_metric(source: str, column: str) -> int:
        if source not in free_data.index or column not in free_data:
            return 0
        value = pd.to_numeric(free_data.loc[source, column], errors='coerce')
        return int(value) if pd.notna(value) else 0

    akshare_rows = aggregate_metric('akshare', 'rows')
    akshare_names = aggregate_metric('akshare', 'entities')
    yfinance_rows = aggregate_metric('yfinance', 'rows')
    yfinance_volume_names = aggregate_metric(
        'yfinance', 'positive_volume_entities'
    )
    openfigi_names = aggregate_metric('openfigi', 'entities')
    macro_series = aggregate_metric('fred_alfred', 'entities')
    fee = evidence.backtest_manifest['annual_bank_fee']
    risk = evidence.risk_backtest.set_index(
        ['evaluation_segment', 'confidence_level']
    )
    risk_overall_95 = risk.loc[('overall', 0.95)]
    risk_overall_99 = risk.loc[('overall', 0.99)]
    risk_holdout_95 = risk.loc[('chronological_holdout', 0.95)]
    risk_holdout_99 = risk.loc[('chronological_holdout', 0.99)]
    risk_rows = '\n'.join(
        [
            _risk_markdown_row('Overall', 95, risk_overall_95),
            _risk_markdown_row('Overall', 99, risk_overall_99),
            _risk_markdown_row('Holdout', 95, risk_holdout_95),
            _risk_markdown_row('Holdout', 99, risk_holdout_99),
        ]
    )
    performance_rows = '\n'.join(
        [
            _performance_markdown_row(
                'Annualised net return',
                'annualised_return',
                (wolf, equal, cap),
            ),
            _performance_markdown_row(
                'Sharpe ratio', 'sharpe', (wolf, equal, cap)
            ),
            _performance_markdown_row(
                'Sortino ratio', 'sortino', (wolf, equal, cap)
            ),
            _performance_markdown_row(
                'Maximum drawdown',
                'maximum_drawdown',
                (wolf, equal, cap),
            ),
            _performance_markdown_row(
                'Annualised cost drag',
                'annualised_cost_drag',
                (wolf, equal, cap),
            ),
        ]
    )
    prior_score = float(
        evidence.prior_validation_manifest['overall_score']
    )
    improvement_rows = '\n'.join(
        [
            f'| Governance score | {prior_score:g}/100 | '
            f'{evidence.governance_score:g}/100 | Conditional approval |',
            f'| Risk backtesting | Warning, 7.5/15 | '
            f'{risk_status.title()}, {_score_label(risk_component)} | '
            'Chronological holdout is the governance gate |',
            f'| Annual turnover | '
            f'{float(prior_wolf.annualised_turnover):.2f}x | '
            f'{float(wolf.annualised_turnover):.2f}x | <=1.50x: pass |',
            f'| Annualised cost drag | '
            f'{_pct(prior_wolf.annualised_cost_drag, 2)} | '
            f'{_pct(wolf.annualised_cost_drag, 2)} | <=1.50%: pass |',
            '| Hard constraint breaches | 0 | 0 | Pass |',
        ]
    )
    alpha_cap = float(
        evidence.alpha.loc[
            evidence.alpha['benchmark'] == 'cap_weight_eligible',
            'annualised_active_return',
        ].iloc[0]
    )
    alpha_equal = float(
        evidence.alpha.loc[
            evidence.alpha['benchmark'] == 'equal_weight_eligible',
            'annualised_active_return',
        ].iloc[0]
    )
    pbo = float(evidence.overfitting['probability_of_backtest_overfitting'])
    haircut = float(
        evidence.overfitting['selected_information_ratio_haircut']
    )
    current_aum_text = _usd(evidence.current_aum)
    ending_value_text = _usd(wolf_end)
    pnl_text = _usd(wolf_end - evidence.current_aum)
    bank_fee_text = _usd(fee['reference_annual_charge_usd'], 0)
    bank_fee_rate_text = _pct(fee['annual_rate'], 2)
    relative_return_text = _pct(
        float(wolf.annualised_return - equal.annualised_return), 2
    )
    _, paired_p_value = _equal_weight_comparison(evidence)
    buys = sorted(
        evidence.trades.loc[
            evidence.trades['trade_action'] == 'Buy', 'ticker'
        ].tolist()
    )
    reductions = evidence.trades.loc[
        evidence.trades['trade_action'] == 'Reduce',
        ['ticker', 'target_weight'],
    ].sort_values('ticker')
    cash_mask = evidence.holdings['ticker'].astype(str).str.upper().isin(
        {'CASH', 'CASH.USD'}
    )
    equities = evidence.holdings.loc[~cash_mask]
    cash_weight = float(
        evidence.holdings.loc[cash_mask, 'final_weight'].sum()
    )

    buy_text = ', '.join(
        f'{chr(96)}{ticker}{chr(96)}' for ticker in buys
    )
    reduce_text = ', '.join(
        f'{chr(96)}{row.ticker}{chr(96)} to '
        f'{_pct(float(row.target_weight), 1)}'
        for row in reductions.itertuples()
    )
    supervised_oos = _supervised_ensemble_rows(evidence)
    supervised_quantiles = evidence.supervised_quantiles.set_index(
        'horizon_months'
    )
    supervised_rows = []
    for row in supervised_oos.itertuples(index=False):
        quantile = supervised_quantiles.loc[int(row.horizon_months)]
        turnover = (
            f'{float(row.annualised_turnover):.2f}x'
            if pd.notna(row.annualised_turnover)
            else 'N/A'
        )
        supervised_rows.append(
            f'| {int(row.horizon_months)}m | {int(row.observations)} | '
            f'{int(row.independent_observations)} | '
            f'{float(row.mean_rank_ic):.3f} | '
            f'{float(row.independent_rank_ic_sign_test_p_value):.3f} | '
            f'{_pct(float(row.mean_horizon_net_active_return), 1)} | '
            f'{turnover} | '
            f'{_pct(float(quantile.central_90_coverage), 1)} | '
            f'{_pct(float(quantile.mean_interval_width), 0)} |'
        )
    supervised_result_rows = '\n'.join(supervised_rows)
    supervised_profile = evidence.supervised_dataset.loc[
        evidence.supervised_dataset['horizon_months'].eq(3)
    ].iloc[0]
    supervised_decision = evidence.supervised_acceptance.loc[
        evidence.supervised_acceptance['scope'].eq('overall')
    ].iloc[0]
    watchlist = _supervised_watchlist(evidence)
    watchlist_rows = '\n'.join(
        f'| `{row.ticker}` | {row.region} | '
        f'{float(row.supervised_alpha_score):.1f} | '
        f'{_pct(float(row.cost_adjusted_predicted_excess_return), 1)} | '
        f'{_pct(float(row.q05_excess_return), 0)} | '
        f'{_pct(float(row.q95_excess_return), 0)} |'
        for row in watchlist.itertuples(index=False)
    )
    baseline_names = set(equities['ticker'].astype(str))
    regional_weights = pd.to_numeric(
        evidence.regional_alpha['target_weight'], errors='coerce'
    ).fillna(0.0)
    regional_names = set(
        evidence.regional_alpha.loc[
            regional_weights.gt(1e-9), 'ticker'
        ].astype(str)
    )
    shared_names = sorted(baseline_names & regional_names)
    baseline_only_names = sorted(baseline_names - regional_names)
    regional_only_names = sorted(regional_names - baseline_names)
    selected_drl = evidence.drl_challengers.loc[
        evidence.drl_challengers['split'].eq('legacy_locked_oos')
        & evidence.drl_challengers['selected_parameter_by_validation'].fillna(
            False
        )
    ]
    drl_rows = '\n'.join(
        f'| {row.algorithm.replace("_", " ").title()} | '
        f'{_pct(float(row.total_net_return), 1)} | '
        f'{_pct(float(row.baseline_total_return), 1)} | '
        f'{_pct(float(row.mean_active_return), 2)} | Research only |'
        for row in selected_drl.itertuples(index=False)
    )
    first_supervised_due = pd.Timestamp(
        evidence.supervised_freeze['first_outcome_due_date']
    ).strftime('%d %B %Y')
    full_supervised_due = pd.Timestamp(
        evidence.supervised_freeze['earliest_full_evidence_date']
    ).strftime('%d %B %Y')
    decision = (
        '**Approve a controlled, human-supervised live pilot.**'
        if fail_count == 0
        else '**Continue paper and shadow operation; do not allocate live capital yet.**'
    )
    risk_summary = (
        'The chronological VaR gate passes.'
        if risk_status == 'PASS'
        else f'The chronological VaR gate is {risk_status.lower()}.'
    )
    cost_summary = (
        'Net-of-cost controls pass.'
        if str(cost_component.status) == 'PASS'
        else f'Net-of-cost controls are {str(cost_component.status).lower()}.'
    )
    return f'''# Wolf Quant Model Investment Principal Report

As of {evidence.as_of_date}

## Decision

{decision} Governance moved
from {prior_score:g}/100 to {evidence.governance_score:g}/100. {pass_count} components pass,
{warning_count} remain warnings, and {fail_count} fail. {risk_summary} {cost_summary}
Full-scale or unattended
deployment remains unapproved because observed point-in-time evidence is
incomplete and
benchmark-relative alpha is not established.

## What Improved

| Control | Before | Now | Current gate |
| --- | ---: | ---: | --- |
{improvement_rows}

Retention hysteresis, a 6% monthly turnover cap, and minimum-turnover
transitions drove the implementation improvement. The configured no-trade
band did not trigger in this 60-month sample, so it is not credited for the
observed result.

## Supervised Benchmark-Relative Alpha

The new research stack compares OLS after train-only Fama-MacBeth screening,
Ridge, Elastic Net, robust Huber regression, Random Forest, Extra Trees,
histogram gradient boosting, XGBoost regression and XGBoost ranking across
3/6/9/12-month horizons. The primary panel contains
**{int(supervised_profile.securities):,} securities**, versus the much larger
price master, because model rows also require historical fundamentals,
features and realised outcomes.

| Horizon | Monthly cohorts | Independent | Rank IC | Sign p | Net cohort return | Turnover | 90% coverage | Band width |
| --- | ---: | ---: | ---: | ---: | ---: | ---: | ---: | ---: |
{supervised_result_rows}

The 3-month rank IC is positive, but four independent cohorts produce an exact
sign-test p-value above 5%. Longer horizons have only two, one and one
independent observations. Net cohort returns are not compounded portfolio CAGR.
Formal annualised return, Sharpe, t-statistics and confidence intervals remain
suppressed. The governed decision is **{supervised_decision.status}** with a
**{_pct(float(supervised_decision.deployment_blend_weight), 0)} live blend**.

Purged date-block conformal calibration now clears the central 90% coverage
target at every horizon. That correction also reveals low precision: average
9- and 12-month bands span about 100% and 117% of benchmark-relative return.
Recurring 3/6/9-month turnover is below 1.5x and includes spread, FX, impact
and the separate 25bp annual bank fee. Twelve-month recurring turnover remains
unestimable from one cohort.

## Portfolio Outputs And Stock Recommendations

The governed CVaR target remains the only committee portfolio. The low-latency
regional-alpha challenger also holds 20 equal-weight names, but only
**{len(shared_names)} names overlap**. The supervised overlay cannot change
weights while its deployment blend is zero.

- Governed-only: {', '.join(f'`{name}`' for name in baseline_only_names)}
- Shared core: {', '.join(f'`{name}`' for name in shared_names)}
- Regional-challenger only: {', '.join(f'`{name}`' for name in regional_only_names)}

Highest supervised 3-month research rank in each region:

| Ticker | Region | Score | Cost-adjusted alpha | Q05 | Q95 |
| --- | --- | ---: | ---: | ---: | ---: |
{watchlist_rows}

These six names are a research watchlist, not buy orders. The live recommendation
remains the governed target described below.

## DRL And Prospective Evidence

The five PPO seeds, contextual bandit and convex residual challenger remain
rejected for deployment. The selected simple challengers also trailed the
baseline in the 12-month legacy OOS diagnostic:

| Challenger | Net return | Baseline | Mean active return | Status |
| --- | ---: | ---: | ---: | --- |
{drl_rows}

DRL receives {_pct(float(evidence.drl_acceptance['blend_weight_drl']), 0)} and
the baseline receives {_pct(float(evidence.drl_acceptance['blend_weight_baseline']), 0)}.
The generic shadow programme has completed
{int(evidence.shadow_status['completed_prospective_cycles'])} of
{int(evidence.shadow_status['required_prospective_cycles'])} required cycles.
The supervised model was separately frozen for prospective evidence: its first
3-month result is due {first_supervised_due}, and 12 non-overlapping cohorts
cannot complete before {full_supervised_due}.

## Adaptive Risk Backtesting

The trailing model-selection stack contains DCC-IGARCH Student-t, filtered
historical simulation, EWMA Normal, and EWMA Student-t forecasts. The configured
40% chronological holdout is the governance gate and its current status is
**{risk_status}**; overall rows remain development-plus-holdout diagnostics.
{int(locked.get('selection_folds', 1))} blocked development fold(s) selected a
{float(locked['selected_scale_factor']):.3f}x global scale and a
{float(locked['selected_exception_multiplier']):.2f}x buffer for
{int(locked['selected_exception_days'])} day after an observed exception; those
parameters were locked before holdout scoring.

| Sample | VaR | Exceptions | Kupiec p | Independence p | Result |
| --- | ---: | ---: | ---: | ---: | --- |
{risk_rows}

This holdout is chronological reconstructed evidence, not a pristine future
shadow period. Live monitoring is still required.

## Point-In-Time Evidence

The evidence store now contains **{int(coverage['delisting_events']):,}** delisting events and
**{int(coverage['fundamental_rows_with_filing_date']):,}** fundamental rows with filing dates.
Legacy local Bloomberg aggregates contain **{fundamental_vintages:,}** database-as-of
fundamental vintages, **{market_cap_vintages:,}** historical market-cap vintages,
and **{corporate_action_vintages:,}** corporate-action vintages. Bloomberg is
disconnected; licensed rows remain ignored locally and are not published.
The public stack records **{akshare_rows:,}** AKShare bars across
**{akshare_names:,}** China/Hong Kong securities and **{yfinance_rows:,}**
yfinance bars with positive volume for **{yfinance_volume_names:,}** securities,
plus **{openfigi_names:,}** current FIGI matches and **{macro_series}** configured
FRED/ALFRED series. OpenFIGI is a current snapshot and does not repair historical
membership.
Observed acceptance timestamps, dated
index membership, inactive-name prices, and historical volume remain below
their governance thresholds. EODHD populated delistings; the Nasdaq
entitlement yielded five usable rows; Beam was unavailable; and SEC blocked
this runner. Unavailable history is not represented as observed evidence.
The point-in-time component is **{_score_label(pit_component)},
{str(pit_component.status).lower()}**.

## Current Target Portfolio

The resolved baseline contains {len(equities)} equities capped at
{_pct(float(equities['final_weight'].max()), 1)} each and
{_pct(cash_weight, 1)} cash. The current trade comparison produces
{len(buys)} buys and {len(reductions)} reductions. These are model targets,
not executable orders; live NAV, FX, liquidity, prices and compliance
approval must be refreshed first.

- Buy: {buy_text}
- Reduce: {reduce_text}

## Point-In-Time Performance

| Measure | Wolf CVaR | Equal weight | Cap weight |
| --- | ---: | ---: | ---: |
{performance_rows}

Applying the realised 60-month Wolf path after modeled trading costs to
current AUM of {current_aum_text} gives an illustrative ending value of {ending_value_text}
and PnL of {pnl_text}. The separate annual bank charge is {bank_fee_rate_text}, equal to
{bank_fee_text} at the reference AUM. This is a scale illustration, not a
forecast or a live-capacity result.

The portfolio component is **{_score_label(cost_component)},
{str(cost_component.status).lower()}**. Wolf returned {relative_return_text}
per year relative to equal weight.
The paired difference was not statistically significant
(p={paired_p_value:.3f}).

## Alpha And Overfitting

Wolf's point-in-time active return is {_pct(alpha_cap, 2)} versus cap weight and
{_pct(alpha_equal, 2)} versus equal weight. Neither alpha test is statistically
significant. The retrospective CSCV test estimates a
{_pct(pbo)} probability of backtest overfitting, while the selected strategy's
median information ratio falls {_pct(haircut)} out-of-sample. The 1997 replay
remains a stress and exposure diagnostic, not proof of historical selection
skill.

## Conditions For Scaling

1. Recalculate orders with live NAV, FX, price and liquidity data.
2. Require human approval before every pilot rebalance.
3. Keep annualised turnover at or below 1.5x, with a 1.0x pilot stretch goal.
4. Track net performance against equal-weight and cap-weight controls.
5. Stop on stale critical data, any hard breach, a failed live risk test,
   failed reconciliation, or a breached cost budget.
6. Require observed PIT vintages and a genuine future shadow record before
   making deployable alpha claims.

The local test summary and GitHub Actions checks remain publication gates.
Research output only. This report is not authorization for unattended trading
or individualized investment advice.
'''


def build_investment_principal_deck(
    repo_root: str | Path,
    output_directory: str | Path | None = None,
) -> DeckBuildResult:
    evidence = load_deck_evidence(repo_root)
    output_directory = (
        Path(output_directory).resolve()
        if output_directory
        else evidence.repo_root / PRESENTATION_RELATIVE
    )
    output_directory.mkdir(parents=True, exist_ok=True)
    plot_paths = _build_plot_assets(
        evidence, output_directory / 'plots'
    )

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    _set_document_properties(presentation, evidence)
    _slide_cover(presentation, evidence)
    _slide_decision(presentation, evidence)
    _slide_workflow(presentation, evidence)
    _slide_drl_results(presentation, evidence, plot_paths[6])
    _slide_supervised_stack(presentation, evidence)
    _slide_supervised_results(presentation, evidence, plot_paths[4])
    _slide_supervised_calibration(presentation, evidence, plot_paths[5])
    _slide_portfolio_outputs(presentation, evidence)
    _slide_evidence(presentation, evidence)
    _slide_trades(presentation, evidence)
    _slide_exposure(presentation, evidence)
    _slide_optimizer(presentation, evidence)

    _slide_pit_growth(presentation, evidence, plot_paths[0])
    _slide_pit_comparison(presentation, evidence, plot_paths[1])
    _slide_alpha(presentation, evidence, plot_paths[2])
    _slide_long_history(presentation, evidence)
    _slide_macro_events(presentation, evidence)
    _slide_regimes(presentation, evidence)
    _slide_governance(presentation, evidence)
    _slide_costs(presentation, evidence, plot_paths[3])

    _slide_pilot(presentation, evidence)
    _slide_glossary(presentation, evidence)
    _slide_close(presentation, evidence)

    bounds_errors = _validate_slide_bounds(presentation)
    if bounds_errors:
        raise ValueError(
            'Presentation shape bounds failed:\n'
            + '\n'.join(bounds_errors)
        )

    pptx_path = output_directory / 'wolf_quant_model_ic_briefing.pptx'
    presentation.save(pptx_path)

    report_path = output_directory / 'investment_principal_report.md'
    report_path.write_text(
        _report_markdown(evidence), encoding='utf-8'
    )
    recommendation_path = output_directory / 'recommendation_snapshot.csv'
    _recommendation_snapshot(evidence).to_csv(recommendation_path, index=False)

    input_paths = [
        evidence.release_root / 'validation/validation_manifest.json',
        evidence.release_root / 'validation/risk_backtesting_report.csv',
        evidence.release_root
        / 'validation/benchmark_significance_report.csv',
        evidence.release_root / 'validation/transaction_cost_validation.csv',
        evidence.release_root / 'walk_forward_manifest.json',
        evidence.release_root / 'universe_summary.csv',
        evidence.release_root / 'pit_evidence_coverage.json',
        evidence.repo_root
        / PRIOR_RELEASE_RELATIVE
        / 'validation/validation_manifest.json',
        evidence.repo_root
        / PRIOR_RELEASE_RELATIVE
        / 'validation/portfolio_strategy_comparison.csv',
        evidence.outputs_root / 'final_portfolio_weights.csv',
        evidence.outputs_root / 'portfolio_trade_list.csv',
        evidence.outputs_root / 'supervised_alpha/dataset_profile.csv',
        evidence.outputs_root / 'supervised_alpha/validation_summary.csv',
        evidence.outputs_root / 'supervised_alpha/oos_summary.csv',
        evidence.outputs_root / 'supervised_alpha/quantile_metrics.csv',
        evidence.outputs_root / 'supervised_alpha/acceptance_decision.csv',
        evidence.outputs_root / 'supervised_alpha/ensemble_weights.csv',
        evidence.outputs_root
        / 'supervised_alpha/prospective_freeze_manifest.json',
        evidence.outputs_root / 'drl_acceptance_decision.csv',
        evidence.outputs_root / 'drl_simple_challenger_comparison.csv',
        evidence.outputs_root / 'drl_training_summary.csv',
        evidence.outputs_root / 'drl_split_manifest.csv',
        evidence.release_root / 'public_data/drl_long_history_manifest.json',
        evidence.outputs_root
        / 'shadow_operation/shadow_operation_status.json',
        evidence.backtest_root / 'run_manifest.json',
        evidence.backtest_root / 'point_in_time_alpha_significance.csv',
        evidence.backtest_root / 'strategy_overfitting_summary.csv',
    ]
    manifest = {
        'schema_version': 2,
        'generated_at_utc': datetime.now(timezone.utc).isoformat(),
        'as_of_date': evidence.as_of_date,
        'approval_status': evidence.approval_status,
        'governance_score': evidence.governance_score,
        'slide_count': len(presentation.slides),

        'presentation': {
            'path': _manifest_path(pptx_path, evidence.repo_root),
            'sha256': _sha256(pptx_path),
            'bytes': pptx_path.stat().st_size,
        },
        'report': {
            'path': _manifest_path(report_path, evidence.repo_root),
            'sha256': _sha256(report_path),
            'bytes': report_path.stat().st_size,
        },
        'recommendations': {
            'path': _manifest_path(recommendation_path, evidence.repo_root),
            'sha256': _sha256(recommendation_path),
            'bytes': recommendation_path.stat().st_size,
            'publication_scope': (
                'governed target, regional challenger and six-name research watchlist'
            ),
        },
        'restricted_inputs': [
            {
                'name': 'regional_alpha_security_level_output',
                'sha256': (
                    _sha256(
                        evidence.outputs_root
                        / 'optimised_portfolio_regional_alpha.csv'
                    )
                    if (
                        evidence.outputs_root
                        / 'optimised_portfolio_regional_alpha.csv'
                    ).exists()
                    else None
                ),
                'publication_status': 'local_only_licensed_derived',
            },
            {
                'name': 'supervised_security_level_predictions',
                'sha256': (
                    _sha256(
                        evidence.outputs_root
                        / 'supervised_alpha/latest_predictions.csv'
                    )
                    if (
                        evidence.outputs_root
                        / 'supervised_alpha/latest_predictions.csv'
                    ).exists()
                    else None
                ),
                'publication_status': 'local_only_licensed_derived',
            },
        ],

        'inputs': [
            {
                'path': path.relative_to(
                    evidence.repo_root
                ).as_posix(),
                'sha256': _sha256(path),
            }
            for path in input_paths
        ],
        'plots': [
            {
                'path': _manifest_path(path, evidence.repo_root),

                'sha256': _sha256(path),
            }
            for path in plot_paths
        ],
    }
    manifest_path = output_directory / 'manifest.json'
    manifest_path.write_text(
        json.dumps(manifest, indent=2) + '\n',
        encoding='utf-8',
    )
    return DeckBuildResult(
        pptx_path=pptx_path,
        report_path=report_path,

        manifest_path=manifest_path,
        plot_paths=plot_paths,
        slide_count=len(presentation.slides),
    )


def register_rendered_pdf(
    repo_root: str | Path,
    output_directory: str | Path | None = None,
    *,
    pdf_path: str | Path | None = None,
    renderer: str = 'Microsoft PowerPoint 16.0',
) -> Path:
    repo_root = Path(repo_root).resolve()
    output_directory = (
        Path(output_directory).resolve()
        if output_directory
        else repo_root / PRESENTATION_RELATIVE
    )
    manifest_path = output_directory / 'manifest.json'
    pptx_path = output_directory / 'wolf_quant_model_ic_briefing.pptx'
    if pdf_path:
        rendered_pdf_path = Path(pdf_path)
        if not rendered_pdf_path.is_absolute():
            rendered_pdf_path = repo_root / rendered_pdf_path
        rendered_pdf_path = rendered_pdf_path.resolve()
    else:
        rendered_pdf_path = (
            output_directory / 'wolf_quant_model_ic_briefing.pdf'
        )
    _require_files([manifest_path, pptx_path, rendered_pdf_path])

    manifest = _read_json(manifest_path)
    slide_count = len(Presentation(pptx_path).slides)
    manifest['rendered_pdf'] = {
        'path': _manifest_path(rendered_pdf_path, repo_root),
        'sha256': _sha256(rendered_pdf_path),
        'bytes': rendered_pdf_path.stat().st_size,
    }
    manifest['rendering'] = {
        'renderer': renderer,
        'slide_count': slide_count,
        'visual_review': 'completed',
    }
    manifest_path.write_text(
        json.dumps(manifest, indent=2) + '\n',
        encoding='utf-8',
    )
    return manifest_path

import json
from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation

from src.reporting.investment_principal_deck import (
    build_investment_principal_deck,
    load_deck_evidence,
    register_rendered_pdf,
)


REPO_ROOT = Path(__file__).resolve().parents[1]


def _presentation_text(presentation: Presentation) -> str:
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                chunks.append(shape.text)
            if getattr(shape, 'has_table', False):
                chunks.extend(
                    cell.text
                    for row in shape.table.rows
                    for cell in row.cells
                )
    return '\n'.join(chunks)


def test_deck_evidence_uses_resolved_portfolio() -> None:
    evidence = load_deck_evidence(REPO_ROOT)

    cash = evidence.holdings.loc[
        evidence.holdings['ticker'].isin(['CASH', 'CASH.USD'])
    ]
    equities = evidence.holdings.drop(index=cash.index)

    assert evidence.holdings['final_weight'].sum() == pytest.approx(1.0)
    assert 0.0 <= cash['final_weight'].sum() <= 0.25 + 1.0e-12
    assert not equities.empty
    assert equities['final_weight'].max() <= 0.05 + 1.0e-12
    assert set(evidence.trades['ticker']) == set(evidence.holdings['ticker'])
    assert (evidence.trades['trade_action'] == 'Buy').any()
    assert evidence.approval_status == 'CONDITIONALLY_APPROVED'
    assert evidence.governance_score == pytest.approx(
        evidence.scorecard['score'].sum()
    )
    assert 0.0 <= evidence.governance_score <= 100.0
    assert evidence.pit_coverage['coverage']['delisting_events'] > 0
    production = evidence.production_pit.set_index('dataset')
    assert int(production.loc['fundamental_vintages', 'rows']) > 0
    risk_gate = evidence.risk_backtest.loc[
        evidence.risk_backtest['governance_gate'].astype(bool)
    ]
    assert not risk_gate.empty
    assert set(risk_gate['status']) == {'PASS'}
    supervised = evidence.supervised_acceptance.set_index('scope').loc[
        'overall'
    ]
    assert supervised['status'] == 'INSUFFICIENT_EVIDENCE'
    assert float(supervised['deployment_blend_weight']) == 0.0
    assert int(
        evidence.supervised_dataset.set_index('horizon_months').loc[3, 'securities']
    ) == 1374
    assert (
        evidence.supervised_quantiles['central_90_coverage'] >= 0.90
    ).all()
    assert evidence.supervised_freeze['legacy_oos_eligible_for_deployment'] is False
    assert not bool(evidence.drl_acceptance['accepted'])
    wolf = evidence.pit_summary.set_index('strategy').loc['wolf_cvar']
    assert wolf['annualised_turnover'] <= 1.5
    assert wolf['annualised_cost_drag'] <= 0.015


def test_deck_builds_with_expected_sections(tmp_path: Path) -> None:
    evidence = load_deck_evidence(REPO_ROOT)
    result = build_investment_principal_deck(
        REPO_ROOT, tmp_path
    )
    presentation = Presentation(result.pptx_path)
    text = _presentation_text(presentation)

    assert result.slide_count == 25
    assert len(presentation.slides) == 25
    fail_count = int(evidence.scorecard['status'].eq('FAIL').sum())
    expected_recommendation = (
        'Approve a controlled live pilot'
        if fail_count == 0
        else 'Continue paper and shadow operation'
    )
    assert expected_recommendation in text
    assert 'Alpha and overfitting' in text
    assert 'Equities to establish' in text
    assert 'Risk holdout passes' in text
    assert 'DRL learned safely, but did not earn capital' in text
    assert 'Design architecture: ML remains inside the governed stack' in text
    assert 'The new supervised alpha research stack' in text
    assert 'How the supervised models compared' in text
    assert 'Supervised signal: encouraging, not yet proven' in text
    assert 'Uncertainty and implementation are now controlled' in text
    assert 'Recommendations: target portfolio versus research challengers' in text
    assert 'NA9.XETRA' in text
    assert '0%' in text
    wolf = evidence.pit_summary.set_index('strategy').loc['wolf_cvar']
    paired_p_value = float(
        evidence.benchmark_significance.loc[
            evidence.benchmark_significance['strategy'].eq('wolf_cvar'),
            'p_value',
        ].iloc[0]
    )
    delistings = int(evidence.pit_coverage['coverage']['delisting_events'])
    assert f'{delistings:,}' in text
    assert 'Free-data checkpoint and publication boundary' in text
    assert f'{float(wolf.annualised_turnover):.2f}x' in text
    assert f'{float(wolf.annualised_cost_drag):.2%}' in text
    assert f'p={paired_p_value:.3f}' in text
    assert f'{evidence.governance_score:g} / 100' in text
    assert result.report_path.exists()
    assert result.manifest_path.exists()
    assert all(path.exists() for path in result.plot_paths)
    assert (tmp_path / 'plots/supervised_model_comparison.png').exists()
    recommendations = pd.read_csv(tmp_path / 'recommendation_snapshot.csv')
    regional_count = int(
        pd.to_numeric(
            evidence.regional_alpha['target_weight'], errors='coerce'
        ).fillna(0.0).gt(0.0).sum()
    )
    assert len(recommendations) == len(evidence.holdings) + regional_count + 6
    assert set(recommendations['recommendation_class']) == {
        'governed_target',
        'regional_alpha_challenger',
        'supervised_research_watchlist',
    }
    assert (
        recommendations.loc[
            recommendations['recommendation_class'].eq(
                'supervised_research_watchlist'
            ),
            'governance_status',
        ]
        == 'research_only_not_a_buy_order'
    ).all()
    manifest = json.loads(result.manifest_path.read_text(encoding='utf-8'))
    assert manifest['recommendations']['path'].endswith(
        'recommendation_snapshot.csv'
    )
    assert len(manifest['restricted_inputs']) == 2
    report = result.report_path.read_text(encoding='utf-8')
    cash_mask = evidence.holdings['ticker'].isin(['CASH', 'CASH.USD'])
    equity_count = int((~cash_mask).sum())
    trade_counts = evidence.trades['trade_action'].value_counts()
    assert f'{equity_count} equities capped at' in report
    assert (
        f"{int(trade_counts.get('Buy', 0))} buys and "
        f"{int(trade_counts.get('Reduce', 0))} reductions"
    ) in report
    assert f'(p={paired_p_value:.3f})' in report
    assert '## Supervised Benchmark-Relative Alpha' in report
    assert 'Implementation and artifact locations' in report
    assert '`src/models/supervised_alpha.py`' in report
    assert '## Portfolio Outputs And Stock Recommendations' in report
    assert '## DRL And Prospective Evidence' in report
    assert '30 November 2026' in report
    assert '31 August 2029' in report


def test_deck_loads_from_public_recommendation_snapshot(
    monkeypatch, tmp_path: Path
) -> None:
    restricted = {
        (
            REPO_ROOT
            / 'reports/outputs/optimised_portfolio_regional_alpha.csv'
        ).resolve(),
        (
            REPO_ROOT
            / 'reports/outputs/supervised_alpha/latest_predictions.csv'
        ).resolve(),
        (
            REPO_ROOT
            / 'reports/outputs/validation/free_data_evidence_summary.csv'
        ).resolve(),
    }
    original_exists = Path.exists

    def published_exists(path: Path) -> bool:
        if path.resolve() in restricted:
            return False
        return original_exists(path)

    monkeypatch.setattr(Path, 'exists', published_exists)
    evidence = load_deck_evidence(REPO_ROOT)

    assert len(evidence.regional_alpha) == 20
    assert len(evidence.supervised_latest) == 6
    assert set(evidence.supervised_latest['horizon_months']) == {3}
    assert int(evidence.free_data_summary['rows'].sum()) > 0

    build_investment_principal_deck(REPO_ROOT, tmp_path)
    rebuilt = pd.read_csv(tmp_path / 'recommendation_snapshot.csv')
    regional = rebuilt.loc[
        rebuilt['recommendation_class'].eq('regional_alpha_challenger')
    ]
    assert regional['model_score'].notna().all()


def test_register_rendered_pdf_accepts_versioned_path(
    tmp_path: Path,
) -> None:
    result = build_investment_principal_deck(REPO_ROOT, tmp_path)
    pdf_path = tmp_path / 'wolf_quant_model_ic_briefing_versioned.pdf'
    pdf_path.write_bytes(b'%PDF-1.4\n%%EOF\n')

    manifest_path = register_rendered_pdf(
        REPO_ROOT,
        tmp_path,
        pdf_path=pdf_path,
        renderer='test renderer',
    )
    manifest = json.loads(manifest_path.read_text(encoding='utf-8'))

    assert result.slide_count == 25
    assert manifest['rendered_pdf']['path'] == str(pdf_path)
    assert manifest['rendering']['renderer'] == 'test renderer'
    assert manifest['rendering']['visual_review'] == 'completed'
